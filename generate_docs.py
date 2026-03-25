#!/usr/bin/env python3
"""Lemuel 프로젝트 산출물 자동 생성 스크립트"""

import os
from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

BASE = os.path.dirname(os.path.abspath(__file__))
DOCS = os.path.join(BASE, "docs")

# ── 공통 스타일 헬퍼 ──────────────────────────────────────────

PROJECT_NAME = "Lemuel (전자상거래 주문·결제·정산 통합 시스템)"
PROJECT_DATE = "2026-03-25"
PROJECT_VER = "1.0"

def make_dirs():
    for d in ["분석", "설계", "구현", "테스트"]:
        os.makedirs(os.path.join(DOCS, d), exist_ok=True)

# ── DOCX 헬퍼 ─────────────────────────────────────────────────

def new_doc(title, subtitle=""):
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Malgun Gothic'
    style.font.size = Pt(10)
    style.paragraph_format.space_after = Pt(4)

    # 표지
    for _ in range(6):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(title)
    run.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(0x1A, 0x47, 0x8A)

    if subtitle:
        p2 = doc.add_paragraph()
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run2 = p2.add_run(subtitle)
        run2.font.size = Pt(14)
        run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    p3 = doc.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run3 = p3.add_run(f"\n\n프로젝트: {PROJECT_NAME}\n버전: {PROJECT_VER}\n작성일: {PROJECT_DATE}")
    run3.font.size = Pt(11)

    doc.add_page_break()
    return doc

def add_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0x1A, 0x47, 0x8A)
    return h

def add_table(doc, headers, rows):
    table = doc.add_table(rows=1+len(rows), cols=len(headers))
    table.style = 'Light Grid Accent 1'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.size = Pt(9)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.rows[ri+1].cells[ci]
            cell.text = str(val)
            for p in cell.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
    return table

# ── XLSX 헬퍼 ─────────────────────────────────────────────────

HEADER_FILL = PatternFill(start_color="1A478A", end_color="1A478A", fill_type="solid")
HEADER_FONT = Font(bold=True, color="FFFFFF", size=10, name="Malgun Gothic")
CELL_FONT = Font(size=10, name="Malgun Gothic")
THIN_BORDER = Border(
    left=Side(style='thin'), right=Side(style='thin'),
    top=Side(style='thin'), bottom=Side(style='thin')
)

def new_wb():
    return Workbook()

def write_sheet(ws, headers, rows, title=None):
    if title:
        ws.title = title
    for ci, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=ci, value=h)
        cell.font = HEADER_FONT
        cell.fill = HEADER_FILL
        cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        cell.border = THIN_BORDER
    for ri, row in enumerate(rows, 2):
        for ci, val in enumerate(row, 1):
            cell = ws.cell(row=ri, column=ci, value=val)
            cell.font = CELL_FONT
            cell.alignment = Alignment(vertical='center', wrap_text=True)
            cell.border = THIN_BORDER
    # auto-width approximation
    for ci in range(1, len(headers)+1):
        max_len = max(len(str(ws.cell(row=r, column=ci).value or "")) for r in range(1, len(rows)+2))
        ws.column_dimensions[chr(64+ci) if ci <= 26 else 'A'].width = min(max_len + 4, 50)

def auto_width(ws, col_count):
    """Better auto-width for columns beyond Z"""
    from openpyxl.utils import get_column_letter
    for ci in range(1, col_count + 1):
        letter = get_column_letter(ci)
        max_len = 0
        for row in ws.iter_rows(min_col=ci, max_col=ci):
            for cell in row:
                if cell.value:
                    max_len = max(max_len, len(str(cell.value)))
        ws.column_dimensions[letter].width = min(max_len + 4, 50)

# ── PPTX 헬퍼 ─────────────────────────────────────────────────

def new_pptx(title, subtitle=""):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return prs

def add_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = PptxPt(14)
    return slide

# ═══════════════════════════════════════════════════════════════
# 1. 분석 단계
# ═══════════════════════════════════════════════════════════════

def gen_현행시스템_분석서():
    doc = new_doc("현행시스템 분석서", "기존 시스템의 구조/기능/데이터 현황 분석")

    add_heading(doc, "1. 문서 개요")
    doc.add_paragraph("본 문서는 Lemuel 전자상거래 주문·결제·정산 통합 시스템의 현행 시스템을 분석한 결과를 기술한다.")

    add_heading(doc, "2. 현행 시스템 구성")
    add_heading(doc, "2.1 시스템 아키텍처", 2)
    doc.add_paragraph("현행 시스템은 모노레포 구조의 헥사고날(Ports & Adapters) 아키텍처를 채택하고 있으며, "
                       "백엔드(Spring Boot, Java 21)와 프론트엔드(React 18, Vite)로 구성된다.")

    add_heading(doc, "2.2 기술 스택", 2)
    add_table(doc,
        ["구분", "기술", "버전", "비고"],
        [
            ["백엔드 프레임워크", "Spring Boot", "3.x", "Java 21 기반"],
            ["프론트엔드", "React", "18.x", "Vite 빌드, TypeScript"],
            ["데이터베이스", "PostgreSQL", "17", "스키마: opslab"],
            ["검색엔진", "Elasticsearch", "8.x", "Nori 한글 분석기"],
            ["캐시", "Redis", "-", "세션/캐시 관리"],
            ["결제 PG", "Toss Payments", "-", "결제 승인/취소"],
            ["컨테이너", "Docker Compose", "-", "개발/운영 환경"],
            ["CI/CD", "GitHub Actions", "-", "자동 빌드/배포"],
            ["코드 품질", "SonarCloud", "-", "정적 분석"],
            ["보안 스캔", "Snyk", "-", "의존성 취약점"],
        ]
    )

    add_heading(doc, "2.3 도메인 구성", 2)
    add_table(doc,
        ["도메인", "역할", "주요 기능"],
        [
            ["user", "인증/인가", "회원가입, 로그인, JWT 발급, 비밀번호 재설정"],
            ["product", "상품 관리", "상품 CRUD, 이미지 관리, Elasticsearch 검색"],
            ["category", "카테고리", "다계층 이커머스 카테고리 관리"],
            ["order", "주문 관리", "주문 생성, 상태 변경, 취소, 환불"],
            ["payment", "결제 관리", "Toss PG 연동, 결제 승인/매입/취소, 환불"],
            ["settlement", "정산 관리", "정산 자동화, 배치 처리, ES 검색"],
            ["coupon", "쿠폰 관리", "쿠폰 발급, 사용, 만료 처리"],
            ["review", "상품 리뷰", "리뷰 등록/조회"],
            ["game", "게임", "오목/바둑 (미완성)"],
            ["common", "공통 모듈", "JWT, Security, Batch, Cache, 예외 처리"],
        ]
    )

    add_heading(doc, "3. 데이터 현황")
    add_heading(doc, "3.1 데이터베이스 스키마", 2)
    doc.add_paragraph("PostgreSQL 17 기반, Flyway 마이그레이션 V1~V21을 통해 스키마가 관리된다.")
    add_table(doc,
        ["테이블명", "설명", "마이그레이션"],
        [
            ["users", "사용자 정보", "V1"],
            ["orders", "주문 정보", "V2"],
            ["payments", "결제 정보", "V2"],
            ["settlements", "정산 정보", "V2"],
            ["products", "상품 정보", "V10"],
            ["password_reset_tokens", "비밀번호 재설정 토큰", "V11"],
            ["categories / tags", "카테고리/태그", "V12"],
            ["ecommerce_categories", "이커머스 카테고리", "V13"],
            ["product_images", "상품 이미지", "V14"],
            ["reviews", "상품 리뷰", "V19"],
            ["coupons", "쿠폰", "V20"],
        ]
    )

    add_heading(doc, "4. 현행 시스템 이슈 및 개선 필요사항")
    issues = [
        "테스트 커버리지 부족: 소스 294개 대비 테스트 17개 (약 5.8%)",
        "게임 도메인(오목/바둑) 미완성 상태 — Controller만 존재",
        "컨트롤러 테스트 부재 (15개 중 14개 미작성)",
        "API 문서화 미비 (Swagger/OpenAPI 미적용)",
        "모니터링/알림 체계 미구축",
    ]
    for issue in issues:
        doc.add_paragraph(issue, style='List Bullet')

    add_heading(doc, "5. 프론트엔드 현황")
    add_table(doc,
        ["화면", "파일", "설명"],
        [
            ["시작 페이지", "StartPage.tsx", "메인 랜딩 페이지"],
            ["로그인", "Login.tsx", "사용자 로그인"],
            ["회원가입", "Register.tsx", "신규 회원가입"],
            ["상품 목록/상세", "ProductPage.tsx", "상품 조회 및 상세"],
            ["장바구니", "CartPage.tsx", "장바구니 관리"],
            ["주문", "OrderPage.tsx", "주문 처리"],
            ["결제 성공/실패", "TossPaymentSuccess/Fail.tsx", "Toss 결제 결과"],
            ["마이페이지", "MyPage.tsx", "사용자 정보/주문내역"],
            ["관리자 대시보드", "AdminDashboardPage.tsx", "관리자 메인"],
            ["정산 대시보드", "SettlementDashboard.tsx", "정산 현황 조회"],
            ["카테고리 관리", "CategoryManagementPage.tsx", "카테고리 CRUD"],
            ["태그 관리", "TagManagementPage.tsx", "태그 CRUD"],
        ]
    )

    doc.save(os.path.join(DOCS, "분석", "현행시스템 분석서.docx"))
    print("  [OK] 현행시스템 분석서.docx")


def gen_요구사항_정의서():
    doc = new_doc("요구사항 정의서 (SRS)", "기능/비기능 요구사항 상세 정의")

    add_heading(doc, "1. 문서 개요")
    doc.add_paragraph("본 문서는 Lemuel 시스템의 기능 및 비기능 요구사항을 정의한다. "
                       "소프트웨어 요구사항 명세서(SRS) 표준을 따른다.")

    add_heading(doc, "2. 기능 요구사항")
    func_reqs = [
        ["FR-USER-001", "회원가입", "사용자는 이메일/비밀번호로 회원가입할 수 있다", "필수"],
        ["FR-USER-002", "로그인", "사용자는 이메일/비밀번호로 로그인하고 JWT 토큰을 발급받는다", "필수"],
        ["FR-USER-003", "비밀번호 재설정", "사용자는 이메일을 통해 비밀번호를 재설정할 수 있다", "필수"],
        ["FR-PROD-001", "상품 등록", "관리자는 상품을 등록할 수 있다 (이미지, 카테고리, 태그 포함)", "필수"],
        ["FR-PROD-002", "상품 검색", "사용자는 키워드로 상품을 검색할 수 있다 (Elasticsearch)", "필수"],
        ["FR-PROD-003", "상품 상태 관리", "관리자는 상품 상태(활성/비활성/품절)를 변경할 수 있다", "필수"],
        ["FR-ORD-001", "주문 생성", "사용자는 상품을 선택하여 주문을 생성할 수 있다", "필수"],
        ["FR-ORD-002", "주문 취소", "사용자는 결제 전 주문을 취소할 수 있다", "필수"],
        ["FR-ORD-003", "주문 상태 조회", "사용자는 주문의 현재 상태를 조회할 수 있다", "필수"],
        ["FR-PAY-001", "결제 승인", "Toss PG를 통해 결제를 승인할 수 있다", "필수"],
        ["FR-PAY-002", "결제 취소/환불", "결제 취소 및 환불을 처리할 수 있다 (멱등성 보장)", "필수"],
        ["FR-STL-001", "일일 정산 생성", "매일 02:00 전날 CAPTURED 결제를 PENDING 정산으로 생성한다", "필수"],
        ["FR-STL-002", "정산 확정", "매일 03:00 PENDING 정산을 CONFIRMED로 확정한다", "필수"],
        ["FR-STL-003", "정산 검색", "Elasticsearch를 통해 정산 내역을 검색할 수 있다", "필수"],
        ["FR-STL-004", "정산 PDF 생성", "정산 내역을 PDF로 다운로드할 수 있다", "필수"],
        ["FR-CAT-001", "카테고리 관리", "다계층 이커머스 카테고리를 CRUD할 수 있다", "필수"],
        ["FR-CPN-001", "쿠폰 발급", "관리자는 쿠폰을 생성/발급할 수 있다", "필수"],
        ["FR-CPN-002", "쿠폰 적용", "사용자는 주문 시 쿠폰을 적용할 수 있다", "필수"],
        ["FR-REV-001", "리뷰 작성", "사용자는 구매한 상품에 리뷰를 작성할 수 있다", "필수"],
    ]
    add_table(doc, ["요구사항 ID", "요구사항명", "설명", "우선순위"], func_reqs)

    add_heading(doc, "3. 비기능 요구사항")
    nfr = [
        ["NFR-PERF-001", "응답 시간", "API 응답 시간 95퍼센타일 500ms 이내", "필수"],
        ["NFR-PERF-002", "동시 처리", "동시 사용자 500명 이상 지원", "필수"],
        ["NFR-SEC-001", "인증", "JWT 기반 인증/인가, 토큰 만료 관리", "필수"],
        ["NFR-SEC-002", "데이터 암호화", "비밀번호 BCrypt 해싱, HTTPS 통신", "필수"],
        ["NFR-SEC-003", "환불 멱등성", "Idempotency-Key 헤더 기반 환불 중복 방지", "필수"],
        ["NFR-SEC-004", "동시성 제어", "환불 시 PESSIMISTIC_WRITE 락 적용", "필수"],
        ["NFR-REL-001", "가용성", "시스템 가용성 99.9% 이상", "권장"],
        ["NFR-REL-002", "배치 중복 방지", "Spring Batch JobRepository 기반 중복 실행 방지", "필수"],
        ["NFR-MNT-001", "테스트 커버리지", "JaCoCo 기준 70% 이상 유지", "필수"],
        ["NFR-MNT-002", "코드 품질", "SonarCloud 정적 분석 통과", "필수"],
        ["NFR-MNT-003", "보안 스캔", "Snyk 의존성 취약점 스캔 통과", "필수"],
    ]
    add_table(doc, ["요구사항 ID", "요구사항명", "설명", "우선순위"], nfr)

    add_heading(doc, "4. 제약사항")
    constraints = [
        "라이선스: AGPL-3.0 (iText 8 의존성)",
        "데이터베이스: PostgreSQL 17 필수 (opslab 스키마)",
        "결제: Toss Payments PG만 지원",
        "검색: Elasticsearch 8.x + Nori 분석기 필수",
    ]
    for c in constraints:
        doc.add_paragraph(c, style='List Bullet')

    doc.save(os.path.join(DOCS, "분석", "요구사항 정의서(SRS).docx"))
    print("  [OK] 요구사항 정의서(SRS).docx")


def gen_RTM():
    wb = new_wb()
    ws = wb.active
    headers = ["요구사항 ID", "요구사항명", "설계 문서", "구현 클래스", "테스트 케이스", "상태"]
    rows = [
        ["FR-USER-001", "회원가입", "시스템 아키텍처 설계서 3.1", "CreateUserService", "CreateUserServiceTest", "구현완료"],
        ["FR-USER-002", "로그인", "시스템 아키텍처 설계서 3.1", "LoginService", "LoginServiceTest", "구현완료"],
        ["FR-USER-003", "비밀번호 재설정", "시스템 아키텍처 설계서 3.1", "PasswordResetService", "-", "테스트 미작성"],
        ["FR-PROD-001", "상품 등록", "시스템 아키텍처 설계서 3.2", "CreateProductService", "-", "테스트 미작성"],
        ["FR-PROD-002", "상품 검색", "시스템 아키텍처 설계서 3.2", "GetProductService", "-", "테스트 미작성"],
        ["FR-PROD-003", "상품 상태 관리", "시스템 아키텍처 설계서 3.2", "ManageProductStatusService", "-", "테스트 미작성"],
        ["FR-ORD-001", "주문 생성", "시스템 아키텍처 설계서 3.3", "CreateOrderService", "CreateOrderServiceTest", "구현완료"],
        ["FR-ORD-002", "주문 취소", "시스템 아키텍처 설계서 3.3", "ChangeOrderStatusService", "-", "테스트 미작성"],
        ["FR-ORD-003", "주문 상태 조회", "시스템 아키텍처 설계서 3.3", "GetOrderService", "-", "테스트 미작성"],
        ["FR-PAY-001", "결제 승인", "시스템 아키텍처 설계서 3.4", "PaymentController", "-", "테스트 미작성"],
        ["FR-PAY-002", "결제 취소/환불", "시스템 아키텍처 설계서 3.4", "PaymentController", "-", "테스트 미작성"],
        ["FR-STL-001", "일일 정산 생성", "시스템 아키텍처 설계서 3.5", "CreateDailySettlementsService", "CreateDailySettlementsServiceTest", "구현완료"],
        ["FR-STL-002", "정산 확정", "시스템 아키텍처 설계서 3.5", "ConfirmDailySettlementsService", "ConfirmDailySettlementsServiceTest", "구현완료"],
        ["FR-STL-003", "정산 검색", "시스템 아키텍처 설계서 3.5", "SettlementSearchController", "SettlementSearchControllerTest", "구현완료"],
        ["FR-STL-004", "정산 PDF 생성", "시스템 아키텍처 설계서 3.5", "GenerateSettlementPdfService", "-", "테스트 미작성"],
        ["FR-CAT-001", "카테고리 관리", "시스템 아키텍처 설계서 3.6", "EcommerceCategoryService", "-", "테스트 미작성"],
        ["FR-CPN-001", "쿠폰 발급", "시스템 아키텍처 설계서 3.7", "CouponService", "CouponServiceTest", "구현완료"],
        ["FR-CPN-002", "쿠폰 적용", "시스템 아키텍처 설계서 3.7", "CouponService", "CouponServiceTest", "구현완료"],
        ["FR-REV-001", "리뷰 작성", "시스템 아키텍처 설계서 3.8", "ReviewController", "-", "테스트 미작성"],
    ]
    write_sheet(ws, headers, rows, "요구사항 추적 매트릭스")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "분석", "요구사항 추적 매트릭스(RTM).xlsx"))
    print("  [OK] 요구사항 추적 매트릭스(RTM).xlsx")


def gen_업무프로세스_정의서():
    doc = new_doc("업무 프로세스 정의서 (BPD)", "현행/개선 업무 프로세스 흐름 정의")

    add_heading(doc, "1. 문서 개요")
    doc.add_paragraph("본 문서는 Lemuel 시스템의 핵심 업무 프로세스를 정의한다.")

    # 주문-결제-정산 프로세스
    add_heading(doc, "2. 주문-결제-정산 프로세스")
    add_heading(doc, "2.1 주문 프로세스", 2)
    add_table(doc,
        ["단계", "액터", "활동", "입력", "출력", "비고"],
        [
            ["1", "고객", "상품 선택 및 장바구니 담기", "상품 정보", "장바구니 데이터", ""],
            ["2", "고객", "주문 생성 요청", "장바구니, 배송지", "주문 정보 (PENDING)", "재고 확인"],
            ["3", "시스템", "재고 차감 및 주문 확정", "주문 정보", "주문 정보 (CONFIRMED)", ""],
            ["4", "고객", "결제 요청", "주문 ID, 결제 수단", "결제 페이지 (Toss)", ""],
            ["5", "Toss PG", "결제 승인", "결제 정보", "결제 결과", ""],
            ["6", "시스템", "결제 완료 처리", "결제 결과", "결제 정보 (CAPTURED)", "주문 상태 → PAID"],
        ]
    )

    add_heading(doc, "2.2 정산 배치 프로세스", 2)
    add_table(doc,
        ["단계", "시간", "활동", "입력", "출력", "비고"],
        [
            ["1", "매일 02:00", "일일 정산 생성", "전날 CAPTURED 결제", "PENDING 정산", "Spring Batch"],
            ["2", "매일 03:00", "정산 확정", "PENDING 정산", "CONFIRMED 정산", "자동 처리"],
            ["3", "매일 03:10", "정산 조정 확정", "PENDING 조정 정산", "CONFIRMED 조정 정산", "환불 반영"],
        ]
    )

    add_heading(doc, "2.3 환불 프로세스", 2)
    add_table(doc,
        ["단계", "액터", "활동", "입력", "출력", "비고"],
        [
            ["1", "고객", "환불 요청", "주문 ID", "환불 요청", ""],
            ["2", "시스템", "환불 처리 (PG 취소)", "결제 정보", "환불 결과", "PESSIMISTIC_WRITE 락"],
            ["3", "시스템", "멱등성 체크", "Idempotency-Key", "중복 방지", "동일 요청 차단"],
            ["4", "시스템", "정산 조정", "환불 정보", "조정 정산 생성", "역정산 처리"],
        ]
    )

    add_heading(doc, "3. 사용자 관리 프로세스")
    add_table(doc,
        ["단계", "액터", "활동", "입력", "출력"],
        [
            ["1", "고객", "회원가입", "이메일, 비밀번호", "사용자 계정 생성"],
            ["2", "고객", "로그인", "이메일, 비밀번호", "JWT 토큰 발급"],
            ["3", "고객", "비밀번호 재설정 요청", "이메일", "재설정 링크 발송"],
            ["4", "고객", "비밀번호 변경", "토큰, 새 비밀번호", "비밀번호 갱신"],
        ]
    )

    add_heading(doc, "4. 상품 관리 프로세스")
    add_table(doc,
        ["단계", "액터", "활동", "입력", "출력"],
        [
            ["1", "관리자", "상품 등록", "상품 정보, 이미지, 카테고리", "상품 생성"],
            ["2", "시스템", "ES 인덱싱", "상품 데이터", "검색 인덱스 갱신"],
            ["3", "관리자", "상품 상태 변경", "상품 ID, 상태", "상태 갱신"],
            ["4", "관리자", "재고 관리", "상품 ID, 수량", "재고 갱신"],
        ]
    )

    doc.save(os.path.join(DOCS, "분석", "업무 프로세스 정의서(BPD).docx"))
    print("  [OK] 업무 프로세스 정의서(BPD).docx")


def gen_유스케이스_명세서():
    doc = new_doc("유스케이스 명세서", "각 유스케이스의 상세 시나리오 기술")

    add_heading(doc, "1. 문서 개요")
    doc.add_paragraph("본 문서는 Lemuel 시스템의 주요 유스케이스를 상세하게 기술한다.")

    usecases = [
        {
            "id": "UC-001", "name": "회원가입",
            "actor": "고객",
            "precondition": "사용자가 시스템에 미가입 상태",
            "main_flow": [
                "1. 사용자가 회원가입 페이지(Register.tsx)에 접근한다.",
                "2. 이메일, 비밀번호, 이름을 입력한다.",
                "3. 시스템은 이메일 중복을 검증한다.",
                "4. 비밀번호를 BCrypt로 해싱하여 저장한다.",
                "5. 사용자 계정이 생성되고 로그인 페이지로 리다이렉트된다.",
            ],
            "alt_flow": "3a. 이메일이 중복인 경우 오류 메시지를 표시한다.",
            "postcondition": "사용자 계정이 DB에 저장된다.",
            "impl": "CreateUserService, AuthController"
        },
        {
            "id": "UC-002", "name": "상품 주문",
            "actor": "고객",
            "precondition": "사용자가 로그인한 상태, 상품이 존재하고 재고가 있음",
            "main_flow": [
                "1. 사용자가 상품을 장바구니에 담는다 (CartContext).",
                "2. 주문 페이지(OrderPage.tsx)에서 주문을 확인한다.",
                "3. 시스템은 재고를 확인하고 주문을 생성한다 (PENDING).",
                "4. Toss 결제 페이지로 리다이렉트된다.",
                "5. 결제 승인 후 시스템이 결제를 확정한다 (CAPTURED).",
                "6. 주문 상태가 PAID로 변경된다.",
            ],
            "alt_flow": "3a. 재고 부족 시 InsufficientStockException 발생. 5a. 결제 실패 시 TossPaymentFail.tsx로 이동.",
            "postcondition": "주문과 결제 레코드가 생성된다.",
            "impl": "CreateOrderService, PaymentController"
        },
        {
            "id": "UC-003", "name": "결제 환불",
            "actor": "고객",
            "precondition": "결제가 CAPTURED 상태",
            "main_flow": [
                "1. 사용자가 마이페이지에서 환불을 요청한다.",
                "2. 시스템은 Idempotency-Key를 확인하여 중복 요청을 방지한다.",
                "3. PESSIMISTIC_WRITE 락을 획득한다.",
                "4. Toss PG에 결제 취소를 요청한다.",
                "5. 환불이 완료되면 정산 조정 레코드를 생성한다.",
            ],
            "alt_flow": "2a. 동일 Idempotency-Key가 존재하면 기존 결과를 반환. 3a. 락 획득 실패 시 재시도.",
            "postcondition": "결제 상태가 REFUNDED, 정산 조정이 생성된다.",
            "impl": "PaymentController, AdjustSettlementForRefundService"
        },
        {
            "id": "UC-004", "name": "일일 정산 처리",
            "actor": "시스템 (배치)",
            "precondition": "전날 CAPTURED 결제 건이 존재",
            "main_flow": [
                "1. 매일 02:00 Spring Batch가 CreateDailySettlements Job을 실행한다.",
                "2. 전날의 CAPTURED 결제를 조회한다.",
                "3. 각 결제에 대해 PENDING 정산 레코드를 생성한다.",
                "4. 03:00 ConfirmDailySettlements Job이 PENDING → CONFIRMED로 확정한다.",
                "5. 03:10 정산 조정 확정 Job이 실행된다.",
            ],
            "alt_flow": "1a. JobRepository에서 중복 실행을 감지하면 Job을 스킵한다.",
            "postcondition": "정산 레코드가 CONFIRMED 상태로 저장된다.",
            "impl": "CreateDailySettlementsService, ConfirmDailySettlementsService"
        },
        {
            "id": "UC-005", "name": "상품 검색",
            "actor": "고객",
            "precondition": "Elasticsearch 인덱스가 구성됨",
            "main_flow": [
                "1. 사용자가 검색어를 입력한다.",
                "2. 시스템은 Elasticsearch Nori 분석기로 한글 토큰화를 수행한다.",
                "3. 검색 결과를 반환한다.",
                "4. 사용자는 카테고리/태그로 필터링할 수 있다.",
            ],
            "alt_flow": "2a. 검색 결과가 없으면 빈 목록을 표시한다.",
            "postcondition": "검색 결과가 화면에 표시된다.",
            "impl": "GetProductService, ProductController"
        },
        {
            "id": "UC-006", "name": "쿠폰 적용",
            "actor": "고객",
            "precondition": "유효한 쿠폰이 존재하고 사용자에게 발급됨",
            "main_flow": [
                "1. 사용자가 주문 시 쿠폰 코드를 입력한다 (CouponInput.tsx).",
                "2. 시스템은 쿠폰 유효성을 검증한다 (만료일, 사용 여부).",
                "3. 할인 금액을 계산하여 주문 금액에서 차감한다.",
                "4. 쿠폰 상태를 사용 완료로 변경한다.",
            ],
            "alt_flow": "2a. 만료된 쿠폰이면 오류 표시. 2b. 이미 사용된 쿠폰이면 오류 표시.",
            "postcondition": "주문 금액이 할인 적용되어 결제된다.",
            "impl": "CouponService, CouponController"
        },
    ]

    for uc in usecases:
        add_heading(doc, f"{uc['id']} — {uc['name']}")
        add_table(doc,
            ["항목", "내용"],
            [
                ["유스케이스 ID", uc["id"]],
                ["유스케이스명", uc["name"]],
                ["액터", uc["actor"]],
                ["사전조건", uc["precondition"]],
                ["후행조건", uc["postcondition"]],
                ["대안흐름", uc["alt_flow"]],
                ["관련 구현", uc["impl"]],
            ]
        )
        add_heading(doc, "기본 흐름", 2)
        for step in uc["main_flow"]:
            doc.add_paragraph(step, style='List Number')
        doc.add_paragraph()

    doc.save(os.path.join(DOCS, "분석", "유스케이스 명세서.docx"))
    print("  [OK] 유스케이스 명세서.docx")


def gen_인터뷰_회의록():
    doc = new_doc("인터뷰 및 회의록", "이해관계자 인터뷰 및 회의 결과 기록")

    add_heading(doc, "1. 킥오프 미팅")
    add_table(doc,
        ["항목", "내용"],
        [
            ["회의일시", "2026-01-06 10:00~12:00"],
            ["참석자", "PM, 백엔드 개발자, 프론트엔드 개발자, 기획자"],
            ["장소", "온라인 (Teams)"],
        ]
    )
    add_heading(doc, "회의 내용", 2)
    items = [
        "프로젝트 범위: 전자상거래 주문·결제·정산 통합 시스템 구축",
        "기술스택 결정: Spring Boot 3.x + React 18 + PostgreSQL 17",
        "아키텍처: 헥사고날 (Ports & Adapters) 패턴 채택",
        "결제 PG: Toss Payments 연동 결정",
        "정산: 일일 배치 처리 방식으로 구현",
        "라이선스: AGPL-3.0 (iText 8 의존성)",
    ]
    for item in items:
        doc.add_paragraph(item, style='List Bullet')

    add_heading(doc, "결정사항", 2)
    decisions = [
        "도메인 모델과 JPA 엔티티를 분리하고 MapStruct로 매핑",
        "환불 처리 시 PESSIMISTIC_WRITE 락 + Idempotency-Key 적용",
        "Elasticsearch Nori 분석기로 한글 검색 지원",
        "CI/CD는 GitHub Actions로 구성",
    ]
    for d in decisions:
        doc.add_paragraph(d, style='List Bullet')

    add_heading(doc, "2. 정산 요구사항 인터뷰")
    add_table(doc,
        ["항목", "내용"],
        [
            ["회의일시", "2026-01-13 14:00~16:00"],
            ["참석자", "PM, 정산 담당 개발자, 재무팀"],
            ["장소", "온라인 (Teams)"],
        ]
    )
    add_heading(doc, "회의 내용", 2)
    items2 = [
        "정산 주기: 일일 정산 (매일 02:00 생성, 03:00 확정)",
        "정산 금액 분할: settlement_amount, fee_amount, net_amount 분리 (V9 마이그레이션)",
        "환불 시 역정산: 별도 조정 레코드 생성",
        "정산 승인 워크플로우: 승인자, 승인일시 필드 추가 (V7 마이그레이션)",
        "ES 검색: 기간/상태/판매자별 정산 검색 기능",
        "PDF 보고서: 정산 내역 PDF 다운로드 기능",
    ]
    for item in items2:
        doc.add_paragraph(item, style='List Bullet')

    add_heading(doc, "3. UI/UX 기획 회의")
    add_table(doc,
        ["항목", "내용"],
        [
            ["회의일시", "2026-01-20 10:00~12:00"],
            ["참석자", "PM, 프론트엔드 개발자, UX 디자이너"],
            ["장소", "온라인 (Teams)"],
        ]
    )
    add_heading(doc, "회의 내용", 2)
    items3 = [
        "고객 화면: 시작페이지 → 상품목록 → 상세 → 장바구니 → 주문 → 결제 흐름",
        "관리자 화면: 대시보드, 상품관리, 카테고리관리, 정산관리",
        "정산 대시보드: 기간별 정산 현황, 차트, 검색 기능",
        "반응형 디자인 적용",
        "Toast 알림 시스템 적용",
    ]
    for item in items3:
        doc.add_paragraph(item, style='List Bullet')

    doc.save(os.path.join(DOCS, "분석", "인터뷰_회의록.docx"))
    print("  [OK] 인터뷰_회의록.docx")


def gen_분석단계_완료보고서():
    prs = new_pptx("분석단계 완료보고서", f"{PROJECT_NAME}\n{PROJECT_DATE}")
    add_slide(prs, "1. 분석 개요", [
        "프로젝트: Lemuel (전자상거래 주문·결제·정산 통합 시스템)",
        "분석 기간: 2026-01-06 ~ 2026-01-31",
        "분석 범위: 사용자, 상품, 주문, 결제, 정산, 쿠폰, 리뷰",
        "분석 방법: 이해관계자 인터뷰, 현행 시스템 분석, 업무 프로세스 분석",
    ])
    add_slide(prs, "2. 주요 산출물", [
        "현행시스템 분석서 — 시스템 구조/기능/데이터 현황",
        "요구사항 정의서(SRS) — 기능 19건, 비기능 11건 정의",
        "요구사항 추적 매트릭스(RTM) — 요구사항-설계-구현-테스트 추적",
        "업무 프로세스 정의서(BPD) — 주문-결제-정산 프로세스 정의",
        "유스케이스 명세서 — 6개 핵심 유스케이스 상세 시나리오",
    ])
    add_slide(prs, "3. 핵심 요구사항 요약", [
        "기능 요구사항: 19건 (사용자 3, 상품 3, 주문 3, 결제 2, 정산 4, 카테고리 1, 쿠폰 2, 리뷰 1)",
        "비기능 요구사항: 11건 (성능 2, 보안 4, 신뢰성 2, 유지보수 3)",
        "핵심 결정: 헥사고날 아키텍처, Toss PG 연동, 일일 배치 정산",
    ])
    add_slide(prs, "4. 리스크 및 이슈", [
        "[리스크] 테스트 커버리지 부족 (현재 약 5.8%, 목표 70%)",
        "[리스크] 게임 도메인 미완성 — 범위 조정 필요",
        "[이슈] API 문서화 미비 — Swagger 적용 필요",
        "[이슈] 모니터링/알림 체계 미구축",
    ])
    add_slide(prs, "5. 다음 단계", [
        "설계 단계 진입 (2026-02-01 ~)",
        "시스템 아키텍처 상세 설계",
        "화면 설계서 (UI/UX) 작성",
        "데이터베이스 상세 설계 (ERD, 테이블 정의서)",
        "API 명세서 작성",
    ])
    prs.save(os.path.join(DOCS, "분석", "분석단계 완료보고서.pptx"))
    print("  [OK] 분석단계 완료보고서.pptx")


# ═══════════════════════════════════════════════════════════════
# 2. 설계 단계
# ═══════════════════════════════════════════════════════════════

def gen_시스템_아키텍처_설계서():
    doc = new_doc("시스템 아키텍처 설계서", "전체 시스템 구조/기술스택/인프라 설계")

    add_heading(doc, "1. 문서 개요")
    doc.add_paragraph("본 문서는 Lemuel 시스템의 전체 아키텍처를 설계한 결과를 기술한다.")

    add_heading(doc, "2. 아키텍처 개요")
    add_heading(doc, "2.1 헥사고날 아키텍처 (Ports & Adapters)", 2)
    doc.add_paragraph(
        "시스템은 헥사고날 아키텍처를 채택하여 도메인 로직의 순수성을 보장한다. "
        "의존성 방향: domain ← application ← adapter. domain 레이어에는 Spring/JPA 의존성이 없다."
    )
    add_table(doc,
        ["레이어", "역할", "포함 요소", "의존성 규칙"],
        [
            ["domain", "순수 비즈니스 로직", "도메인 객체, 값 객체, Enum", "외부 의존성 없음 (POJO)"],
            ["application", "유스케이스 구현", "Service, UseCase, Port 인터페이스", "domain만 의존"],
            ["adapter/in", "외부 → 내부 연결", "Controller, Request/Response DTO", "application 의존"],
            ["adapter/out", "내부 → 외부 연결", "JPA Entity, Repository, Mapper", "application 의존"],
        ]
    )

    add_heading(doc, "2.2 패키지 구조", 2)
    doc.add_paragraph(
        "각 도메인은 다음 디렉토리 구조를 따른다:\n"
        "{domain}/adapter/in/web/ — REST Controller, DTO\n"
        "{domain}/adapter/out/persistence/ — JPA Entity, Repository, Mapper\n"
        "{domain}/adapter/out/search/ — Elasticsearch Adapter\n"
        "{domain}/application/port/in/ — UseCase 인터페이스\n"
        "{domain}/application/port/out/ — Port 인터페이스\n"
        "{domain}/application/service/ — UseCase 구현체\n"
        "{domain}/domain/ — 순수 도메인 객체"
    )

    add_heading(doc, "3. 기술 스택 설계")
    add_table(doc,
        ["구분", "기술", "버전", "용도"],
        [
            ["Language", "Java", "21", "백엔드 메인 언어"],
            ["Framework", "Spring Boot", "3.x", "백엔드 프레임워크"],
            ["Build", "Gradle", "8.x", "빌드 도구"],
            ["Frontend", "React", "18.x", "프론트엔드 프레임워크"],
            ["Frontend Build", "Vite", "5.x", "프론트엔드 빌드"],
            ["Language", "TypeScript", "5.x", "프론트엔드 언어"],
            ["DB", "PostgreSQL", "17", "메인 데이터베이스"],
            ["Search", "Elasticsearch", "8.x", "전문 검색"],
            ["Cache", "Redis", "-", "캐시/세션"],
            ["PG", "Toss Payments", "-", "결제 연동"],
            ["Batch", "Spring Batch", "-", "정산 배치 처리"],
            ["Migration", "Flyway", "-", "DB 마이그레이션"],
            ["Mapping", "MapStruct", "-", "도메인↔엔티티 매핑"],
            ["Container", "Docker Compose", "-", "컨테이너 오케스트레이션"],
            ["CI/CD", "GitHub Actions", "-", "자동화 빌드/배포"],
        ]
    )

    add_heading(doc, "4. 인프라 설계")
    add_heading(doc, "4.1 Docker Compose 구성", 2)
    add_table(doc,
        ["서비스", "이미지", "포트", "용도"],
        [
            ["app", "lemuel-backend", "8080", "백엔드 API 서버"],
            ["frontend", "lemuel-frontend", "5173", "프론트엔드 서버"],
            ["postgres", "postgres:17", "5432", "데이터베이스"],
            ["elasticsearch", "elasticsearch:8", "9200", "검색엔진"],
            ["redis", "redis", "6379", "캐시"],
        ]
    )

    add_heading(doc, "4.2 CI/CD 파이프라인", 2)
    doc.add_paragraph(
        "GitHub Actions (ci.yml):\n"
        "1. 변경 파일 감지 → backend/frontend 선택적 실행\n"
        "2. Backend: Gradle 빌드 → JaCoCo → SonarCloud → Snyk → GHCR push\n"
        "3. Frontend: TypeScript 체크 → ESLint → Vite 빌드 → Snyk → GHCR push"
    )

    add_heading(doc, "5. 보안 설계")
    add_table(doc,
        ["항목", "구현 방식", "비고"],
        [
            ["인증", "JWT (Access Token + Refresh Token)", "JWT_ISSUER, JWT_SECRET 환경변수"],
            ["인가", "Spring Security + Role 기반", "USER, ADMIN 역할"],
            ["비밀번호", "BCrypt 해싱", ""],
            ["환불 동시성", "PESSIMISTIC_WRITE 락", "DB 레벨 배타적 락"],
            ["환불 멱등성", "Idempotency-Key 헤더", "동일 요청 중복 방지"],
            ["배치 중복", "Spring Batch JobRepository", "동일 파라미터 재실행 방지"],
            ["통신", "HTTPS", "TLS 암호화"],
        ]
    )

    doc.save(os.path.join(DOCS, "설계", "시스템 아키텍처 설계서.docx"))
    print("  [OK] 시스템 아키텍처 설계서.docx")


def gen_화면설계서():
    """PDF 화면설계서 - reportlab 사용"""
    filepath = os.path.join(DOCS, "설계", "화면설계서(UI_UX).pdf")
    c = canvas.Canvas(filepath, pagesize=A4)
    w, h = A4

    # Try to register a Korean font
    font_name = "Helvetica"
    try:
        # Try common Windows Korean font paths
        for fpath in [
            "C:/Windows/Fonts/malgun.ttf",
            "C:/Windows/Fonts/NanumGothic.ttf",
        ]:
            if os.path.exists(fpath):
                pdfmetrics.registerFont(TTFont("Korean", fpath))
                font_name = "Korean"
                break
    except Exception:
        pass

    def draw_title_page():
        c.setFont(font_name, 28)
        c.drawCentredString(w/2, h - 250, "화면설계서 (UI/UX)")
        c.setFont(font_name, 14)
        c.drawCentredString(w/2, h - 290, "화면 레이아웃, 네비게이션, 와이어프레임")
        c.setFont(font_name, 11)
        c.drawCentredString(w/2, h - 340, f"프로젝트: {PROJECT_NAME}")
        c.drawCentredString(w/2, h - 360, f"버전: {PROJECT_VER}  |  작성일: {PROJECT_DATE}")
        c.showPage()

    def draw_page(title, items):
        c.setFont(font_name, 18)
        c.drawString(2*cm, h - 2*cm, title)
        c.setFont(font_name, 11)
        y = h - 3.5*cm
        for item in items:
            if y < 3*cm:
                c.showPage()
                c.setFont(font_name, 11)
                y = h - 2*cm
            c.drawString(2.5*cm, y, item)
            y -= 0.7*cm
        c.showPage()

    draw_title_page()

    # 네비게이션 구조
    draw_page("1. 화면 네비게이션 구조", [
        "[고객 영역]",
        "  StartPage → Login/Register",
        "  StartPage → ProductPage (상품 목록/상세)",
        "  ProductPage → CartPage (장바구니)",
        "  CartPage → OrderPage (주문)",
        "  OrderPage → TossPaymentSuccess / TossPaymentFail",
        "  Header → MyPage (마이페이지)",
        "  Login → ForgotPassword → ResetPassword",
        "",
        "[관리자 영역]",
        "  AdminLoginPage → AdminDashboardPage",
        "  AdminDashboardPage → CategoryManagementPage",
        "  AdminDashboardPage → EcommerceCategoryAdmin",
        "  AdminDashboardPage → TagManagementPage",
        "  AdminDashboardPage → SettlementAdmin",
        "  AdminDashboardPage → SettlementDashboard",
    ])

    screens = [
        ("2. StartPage (시작 페이지)", [
            "레이아웃: Layout.tsx 공통 레이아웃 적용",
            "상단: 네비게이션 바 (로고, 검색, 장바구니, 로그인)",
            "메인: 추천 상품, 카테고리별 상품 목록",
            "하단: 푸터 (회사 정보, 링크)",
        ]),
        ("3. Login / Register (로그인/회원가입)", [
            "Login.tsx: 이메일/비밀번호 입력 폼, 로그인 버튼, 회원가입 링크",
            "Register.tsx: 이메일/비밀번호/이름 입력 폼, 가입 버튼",
            "ForgotPassword.tsx: 이메일 입력 → 재설정 링크 발송",
            "ResetPassword.tsx: 새 비밀번호 입력 폼",
        ]),
        ("4. ProductPage (상품)", [
            "상품 목록: Card.tsx 기반 그리드 레이아웃",
            "검색: 키워드 검색 + 카테고리/태그 필터",
            "상품 상세: 이미지, 가격, 설명, 리뷰(ReviewList.tsx), 리뷰 작성(ReviewForm.tsx)",
            "장바구니 담기: CartContext를 통한 상태 관리",
        ]),
        ("5. CartPage / OrderPage (장바구니/주문)", [
            "CartPage.tsx: 상품 목록, 수량 변경, 삭제, 쿠폰 적용(CouponInput.tsx)",
            "OrderPage.tsx: 주문 확인, 배송지 입력, 결제 요청",
            "TossPaymentSuccess.tsx: 결제 성공 후 주문 확인",
            "TossPaymentFail.tsx: 결제 실패 안내 및 재시도",
        ]),
        ("6. AdminDashboardPage (관리자)", [
            "관리자 대시보드: 주요 지표 카드, 최근 주문, 정산 현황",
            "CategoryManagementPage.tsx: 카테고리 트리 CRUD",
            "TagManagementPage.tsx: 태그 관리",
            "상품 관리: CreateProductForm.tsx, ImageUpload.tsx, ProductList.tsx",
        ]),
        ("7. SettlementDashboard (정산)", [
            "SettlementDashboard.tsx: 기간별 정산 현황 대시보드",
            "SettlementDashboardImproved.tsx: 개선된 정산 대시보드",
            "SettlementAdmin.tsx: 정산 관리 (검색, 상태 변경, PDF 다운로드)",
            "DateRangePicker.tsx: 기간 선택 컴포넌트",
            "StatusBadge.tsx: 정산 상태 표시",
        ]),
        ("8. 공통 컴포넌트", [
            "Layout.tsx: 전체 레이아웃 (헤더, 사이드바, 콘텐츠)",
            "Card.tsx: 카드형 UI 컴포넌트",
            "Spinner.tsx: 로딩 스피너",
            "LoadingSkeleton.tsx: 스켈레톤 로딩",
            "Toast.tsx + ToastContext.tsx: 토스트 알림 시스템",
            "EmptyState.tsx: 빈 상태 표시",
            "StarRating.tsx: 별점 표시",
        ]),
    ]
    for title, items in screens:
        draw_page(title, items)

    c.save()
    print("  [OK] 화면설계서(UI_UX).pdf")


def gen_데이터베이스_설계서():
    doc = new_doc("데이터베이스 설계서 (ERD)", "논리/물리 ERD, 테이블 정의서")

    add_heading(doc, "1. 문서 개요")
    doc.add_paragraph("본 문서는 Lemuel 시스템의 데이터베이스 논리/물리 설계를 기술한다. "
                       "PostgreSQL 17 기반, 스키마명 opslab.")

    add_heading(doc, "2. ERD 개요 (엔티티 관계)")
    doc.add_paragraph(
        "users (1) ──< (N) orders (1) ──< (N) payments (1) ──< (N) settlements\n"
        "users (1) ──< (N) reviews\n"
        "products (1) ──< (N) order_items\n"
        "products (1) ──< (N) product_images\n"
        "products (N) >──< (N) tags (via product_tags)\n"
        "products (N) >── (1) categories\n"
        "ecommerce_categories — 자기 참조 (parent_id)\n"
        "coupons — 독립 엔티티"
    )

    add_heading(doc, "3. 테이블 목록")
    add_table(doc,
        ["No", "테이블명", "설명", "마이그레이션", "비고"],
        [
            ["1", "users", "사용자 정보", "V1", "인증/인가"],
            ["2", "orders", "주문 정보", "V2", "주문 관리"],
            ["3", "payments", "결제 정보", "V2", "PG 연동"],
            ["4", "settlements", "정산 정보", "V2, V4, V7, V9", "배치 정산"],
            ["5", "products", "상품 정보", "V10", "상품 관리"],
            ["6", "password_reset_tokens", "비밀번호 재설정 토큰", "V11", ""],
            ["7", "categories", "상품 카테고리", "V12", ""],
            ["8", "tags", "상품 태그", "V12", ""],
            ["9", "product_tags", "상품-태그 매핑", "V12", "M:N"],
            ["10", "ecommerce_categories", "이커머스 카테고리", "V13", "다계층"],
            ["11", "product_images", "상품 이미지", "V14", ""],
            ["12", "reviews", "상품 리뷰", "V19", ""],
            ["13", "coupons", "쿠폰", "V20", ""],
            ["14", "settlement_index_queue", "정산 ES 인덱스 큐", "V5", ""],
            ["15", "settlement_schedule_config", "정산 스케줄 설정", "V6", ""],
        ]
    )

    add_heading(doc, "4. Flyway 마이그레이션 이력")
    add_table(doc,
        ["버전", "파일명", "설명"],
        [
            ["V1", "V1__init.sql", "초기 스키마 (users 등)"],
            ["V2", "V2__create_order_payment_settlement.sql", "주문/결제/정산 테이블"],
            ["V3", "V3__add_indexes_and_constraints.sql", "인덱스 및 제약조건 추가"],
            ["V4", "V4__refunds_and_settlement_adjustments.sql", "환불/정산조정"],
            ["V5", "V5__settlement_index_queue.sql", "정산 ES 인덱스 큐"],
            ["V6", "V6__settlement_schedule_config.sql", "정산 스케줄 설정"],
            ["V7", "V7__add_settlement_approval_fields.sql", "정산 승인 필드"],
            ["V8", "V8__add_user_status_column.sql", "사용자 상태 컬럼"],
            ["V9", "V9__alter_settlements_split_amount.sql", "정산 금액 분할"],
            ["V10", "V10__create_products_table.sql", "상품 테이블"],
            ["V11", "V11__create_password_reset_tokens_table.sql", "비밀번호 재설정 토큰"],
            ["V12", "V12__create_categories_and_tags_tables.sql", "카테고리/태그"],
            ["V13", "V13__create_ecommerce_categories_table.sql", "이커머스 카테고리"],
            ["V14", "V14__create_product_images_table.sql", "상품 이미지"],
            ["V15", "V15__add_product_id_to_orders.sql", "주문에 상품ID 추가"],
            ["V16", "V16__fix_pg_transaction_id_length.sql", "PG 트랜잭션ID 길이 수정"],
            ["V17", "V17__seed_data.sql", "시드 데이터"],
            ["V18", "V18__add_seed_manager.sql", "시드 매니저"],
            ["V19", "V19__create_reviews_table.sql", "리뷰 테이블"],
            ["V20", "V20__create_coupons_table.sql", "쿠폰 테이블"],
            ["V21", "V21__seed_january_2026_data.sql", "2026년 1월 시드"],
        ]
    )

    add_heading(doc, "5. Elasticsearch 인덱스")
    add_table(doc,
        ["인덱스명", "용도", "분석기", "주요 필드"],
        [
            ["settlements", "정산 검색", "Nori (한글)", "정산ID, 금액, 상태, 날짜, 판매자"],
            ["products", "상품 검색", "Nori (한글)", "상품명, 설명, 카테고리, 태그, 가격"],
        ]
    )

    doc.save(os.path.join(DOCS, "설계", "데이터베이스 설계서(ERD).docx"))
    print("  [OK] 데이터베이스 설계서(ERD).docx")


def gen_테이블_정의서():
    wb = new_wb()
    ws = wb.active

    headers = ["테이블명", "컬럼명", "데이터타입", "PK", "FK", "NOT NULL", "DEFAULT", "설명"]
    tables_data = [
        # users
        ["users", "id", "BIGSERIAL", "Y", "", "Y", "", "사용자 ID"],
        ["users", "email", "VARCHAR(255)", "", "", "Y", "", "이메일 (UNIQUE)"],
        ["users", "password", "VARCHAR(255)", "", "", "Y", "", "BCrypt 해싱 비밀번호"],
        ["users", "name", "VARCHAR(100)", "", "", "Y", "", "사용자 이름"],
        ["users", "role", "VARCHAR(20)", "", "", "Y", "USER", "역할 (USER/ADMIN)"],
        ["users", "status", "VARCHAR(20)", "", "", "Y", "ACTIVE", "상태 (V8)"],
        ["users", "created_at", "TIMESTAMP", "", "", "Y", "NOW()", "생성일시"],
        ["users", "updated_at", "TIMESTAMP", "", "", "", "", "수정일시"],
        # orders
        ["orders", "id", "BIGSERIAL", "Y", "", "Y", "", "주문 ID"],
        ["orders", "user_id", "BIGINT", "", "users(id)", "Y", "", "주문자"],
        ["orders", "product_id", "BIGINT", "", "products(id)", "", "", "상품 (V15)"],
        ["orders", "status", "VARCHAR(20)", "", "", "Y", "PENDING", "주문상태"],
        ["orders", "total_amount", "DECIMAL(15,2)", "", "", "Y", "", "총 금액"],
        ["orders", "created_at", "TIMESTAMP", "", "", "Y", "NOW()", "주문일시"],
        ["orders", "updated_at", "TIMESTAMP", "", "", "", "", "수정일시"],
        # payments
        ["payments", "id", "BIGSERIAL", "Y", "", "Y", "", "결제 ID"],
        ["payments", "order_id", "BIGINT", "", "orders(id)", "Y", "", "주문"],
        ["payments", "pg_transaction_id", "VARCHAR(200)", "", "", "", "", "PG 거래번호 (V16)"],
        ["payments", "status", "VARCHAR(20)", "", "", "Y", "", "결제상태"],
        ["payments", "amount", "DECIMAL(15,2)", "", "", "Y", "", "결제금액"],
        ["payments", "payment_method", "VARCHAR(50)", "", "", "", "", "결제수단"],
        ["payments", "captured_at", "TIMESTAMP", "", "", "", "", "매입일시"],
        ["payments", "created_at", "TIMESTAMP", "", "", "Y", "NOW()", "생성일시"],
        # settlements
        ["settlements", "id", "BIGSERIAL", "Y", "", "Y", "", "정산 ID"],
        ["settlements", "payment_id", "BIGINT", "", "payments(id)", "Y", "", "결제"],
        ["settlements", "status", "VARCHAR(20)", "", "", "Y", "PENDING", "정산상태"],
        ["settlements", "settlement_amount", "DECIMAL(15,2)", "", "", "Y", "", "정산금액 (V9)"],
        ["settlements", "fee_amount", "DECIMAL(15,2)", "", "", "", "", "수수료 (V9)"],
        ["settlements", "net_amount", "DECIMAL(15,2)", "", "", "", "", "순수익 (V9)"],
        ["settlements", "approved_by", "VARCHAR(100)", "", "", "", "", "승인자 (V7)"],
        ["settlements", "approved_at", "TIMESTAMP", "", "", "", "", "승인일시 (V7)"],
        ["settlements", "settled_at", "DATE", "", "", "", "", "정산일"],
        ["settlements", "created_at", "TIMESTAMP", "", "", "Y", "NOW()", "생성일시"],
        # products
        ["products", "id", "BIGSERIAL", "Y", "", "Y", "", "상품 ID"],
        ["products", "name", "VARCHAR(255)", "", "", "Y", "", "상품명"],
        ["products", "description", "TEXT", "", "", "", "", "상품 설명"],
        ["products", "price", "DECIMAL(15,2)", "", "", "Y", "", "가격"],
        ["products", "stock", "INTEGER", "", "", "Y", "0", "재고"],
        ["products", "status", "VARCHAR(20)", "", "", "Y", "ACTIVE", "상품상태"],
        ["products", "category_id", "BIGINT", "", "categories(id)", "", "", "카테고리"],
        ["products", "created_at", "TIMESTAMP", "", "", "Y", "NOW()", "생성일시"],
        # product_images
        ["product_images", "id", "BIGSERIAL", "Y", "", "Y", "", "이미지 ID"],
        ["product_images", "product_id", "BIGINT", "", "products(id)", "Y", "", "상품"],
        ["product_images", "url", "VARCHAR(500)", "", "", "Y", "", "이미지 URL"],
        ["product_images", "sort_order", "INTEGER", "", "", "", "0", "정렬순서"],
        # ecommerce_categories
        ["ecommerce_categories", "id", "BIGSERIAL", "Y", "", "Y", "", "카테고리 ID"],
        ["ecommerce_categories", "name", "VARCHAR(100)", "", "", "Y", "", "카테고리명"],
        ["ecommerce_categories", "slug", "VARCHAR(100)", "", "", "Y", "", "슬러그 (UNIQUE)"],
        ["ecommerce_categories", "parent_id", "BIGINT", "", "ecommerce_categories(id)", "", "", "상위카테고리"],
        ["ecommerce_categories", "depth", "INTEGER", "", "", "Y", "0", "깊이"],
        # reviews
        ["reviews", "id", "BIGSERIAL", "Y", "", "Y", "", "리뷰 ID"],
        ["reviews", "product_id", "BIGINT", "", "products(id)", "Y", "", "상품"],
        ["reviews", "user_id", "BIGINT", "", "users(id)", "Y", "", "작성자"],
        ["reviews", "rating", "INTEGER", "", "", "Y", "", "별점 (1-5)"],
        ["reviews", "content", "TEXT", "", "", "", "", "리뷰 내용"],
        ["reviews", "created_at", "TIMESTAMP", "", "", "Y", "NOW()", "작성일시"],
        # coupons
        ["coupons", "id", "BIGSERIAL", "Y", "", "Y", "", "쿠폰 ID"],
        ["coupons", "code", "VARCHAR(50)", "", "", "Y", "", "쿠폰 코드 (UNIQUE)"],
        ["coupons", "type", "VARCHAR(20)", "", "", "Y", "", "쿠폰 유형"],
        ["coupons", "discount_value", "DECIMAL(15,2)", "", "", "Y", "", "할인 값"],
        ["coupons", "expires_at", "TIMESTAMP", "", "", "", "", "만료일시"],
        ["coupons", "used", "BOOLEAN", "", "", "Y", "FALSE", "사용 여부"],
    ]

    write_sheet(ws, headers, tables_data, "테이블 정의서")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "설계", "테이블 정의서.xlsx"))
    print("  [OK] 테이블 정의서.xlsx")


def gen_인터페이스_설계서():
    doc = new_doc("인터페이스 설계서", "내/외부 시스템 간 연동 인터페이스 정의")

    add_heading(doc, "1. 문서 개요")
    doc.add_paragraph("본 문서는 Lemuel 시스템의 내/외부 시스템 간 연동 인터페이스를 정의한다.")

    add_heading(doc, "2. 외부 시스템 인터페이스")
    add_heading(doc, "2.1 Toss Payments PG 연동", 2)
    add_table(doc,
        ["인터페이스 ID", "방향", "프로토콜", "설명", "인증 방식"],
        [
            ["IF-PG-001", "Outbound", "HTTPS REST", "결제 승인 요청", "Secret Key (Basic Auth)"],
            ["IF-PG-002", "Outbound", "HTTPS REST", "결제 취소 요청", "Secret Key (Basic Auth)"],
            ["IF-PG-003", "Inbound", "HTTPS Callback", "결제 결과 콜백", "Webhook Signature"],
        ]
    )
    doc.add_paragraph()
    doc.add_paragraph("Toss Payments 연동 흐름:")
    steps = [
        "1. 클라이언트에서 Toss SDK로 결제 요청",
        "2. Toss 결제 페이지에서 결제 진행",
        "3. 결제 완료 후 successUrl로 리다이렉트 (TossPaymentSuccess.tsx)",
        "4. 백엔드에서 Toss API로 결제 승인 확인",
        "5. 결제 상태를 CAPTURED로 변경",
    ]
    for s in steps:
        doc.add_paragraph(s, style='List Number')

    add_heading(doc, "2.2 이메일 발송 (SMTP)", 2)
    add_table(doc,
        ["인터페이스 ID", "방향", "프로토콜", "설명", "인증 방식"],
        [
            ["IF-MAIL-001", "Outbound", "SMTP", "비밀번호 재설정 이메일 발송", "MAIL_USERNAME/PASSWORD"],
        ]
    )

    add_heading(doc, "3. 내부 시스템 인터페이스")
    add_heading(doc, "3.1 Elasticsearch 연동", 2)
    add_table(doc,
        ["인터페이스 ID", "방향", "프로토콜", "설명"],
        [
            ["IF-ES-001", "Outbound", "HTTP REST", "상품 인덱싱 (색인)"],
            ["IF-ES-002", "Outbound", "HTTP REST", "상품 검색 쿼리"],
            ["IF-ES-003", "Outbound", "HTTP REST", "정산 인덱싱 (IndexSettlementService)"],
            ["IF-ES-004", "Outbound", "HTTP REST", "정산 검색 쿼리 (SettlementSearchController)"],
        ]
    )

    add_heading(doc, "3.2 프론트엔드 ↔ 백엔드 API", 2)
    doc.add_paragraph(
        "프론트엔드는 axios 인스턴스(frontend/src/api/axios.ts)를 통해 백엔드 REST API와 통신한다. "
        "JWT 토큰을 Authorization 헤더에 포함하는 인터셉터가 설정되어 있다."
    )
    add_table(doc,
        ["API 모듈", "파일", "주요 엔드포인트"],
        [
            ["인증", "api/auth.ts", "/api/auth/login, /api/auth/register"],
            ["상품", "api/product.ts", "/api/products"],
            ["주문", "api/order.ts", "/api/orders"],
            ["결제", "api/payment.ts", "/api/payments"],
            ["정산", "api/settlement.ts", "/api/settlements"],
            ["카테고리", "api/category.ts", "/api/categories"],
            ["쿠폰", "api/coupon.ts", "/api/coupons"],
            ["리뷰", "api/review.ts", "/api/reviews"],
            ["태그", "api/tag.ts", "/api/tags"],
            ["환불", "api/refund.ts", "/api/refunds"],
            ["관리자", "api/admin.ts", "/api/admin/*"],
        ]
    )

    doc.save(os.path.join(DOCS, "설계", "인터페이스 설계서.docx"))
    print("  [OK] 인터페이스 설계서.docx")


def gen_API_명세서():
    wb = new_wb()
    ws = wb.active
    headers = ["도메인", "HTTP Method", "Endpoint", "설명", "인증", "Request Body", "Response", "Controller"]
    rows = [
        # Auth
        ["user", "POST", "/api/auth/register", "회원가입", "N", "email, password, name", "User", "AuthController"],
        ["user", "POST", "/api/auth/login", "로그인", "N", "email, password", "JWT Token", "AuthController"],
        ["user", "POST", "/api/auth/forgot-password", "비밀번호 재설정 요청", "N", "email", "Success", "AuthController"],
        ["user", "POST", "/api/auth/reset-password", "비밀번호 재설정", "N", "token, newPassword", "Success", "AuthController"],
        ["user", "GET", "/api/users/{id}", "사용자 조회", "Y", "-", "User", "UserController"],
        ["user", "PUT", "/api/users/{id}", "사용자 수정", "Y", "User 정보", "User", "UserController"],
        # Product
        ["product", "GET", "/api/products", "상품 목록 조회", "N", "-", "Page<Product>", "ProductController"],
        ["product", "GET", "/api/products/{id}", "상품 상세 조회", "N", "-", "Product", "ProductController"],
        ["product", "POST", "/api/products", "상품 등록", "Y(ADMIN)", "Product 정보", "Product", "ProductController"],
        ["product", "PUT", "/api/products/{id}", "상품 수정", "Y(ADMIN)", "Product 정보", "Product", "ProductController"],
        ["product", "DELETE", "/api/products/{id}", "상품 삭제", "Y(ADMIN)", "-", "Success", "ProductController"],
        ["product", "POST", "/api/products/{id}/images", "상품 이미지 업로드", "Y(ADMIN)", "MultipartFile", "ProductImage", "ProductImageController"],
        ["product", "GET", "/api/products/search", "상품 검색 (ES)", "N", "keyword, category", "Page<Product>", "ProductController"],
        # Category
        ["category", "GET", "/api/categories", "카테고리 목록", "N", "-", "List<Category>", "CategoryController"],
        ["category", "POST", "/api/categories", "카테고리 생성", "Y(ADMIN)", "name, slug, parentId", "Category", "AdminEcommerceCategoryController"],
        ["category", "PUT", "/api/categories/{id}", "카테고리 수정", "Y(ADMIN)", "name, slug", "Category", "AdminEcommerceCategoryController"],
        ["category", "DELETE", "/api/categories/{id}", "카테고리 삭제", "Y(ADMIN)", "-", "Success", "AdminEcommerceCategoryController"],
        # Tag
        ["product", "GET", "/api/tags", "태그 목록", "N", "-", "List<Tag>", "TagController"],
        ["product", "POST", "/api/tags", "태그 생성", "Y(ADMIN)", "name", "Tag", "TagController"],
        # Order
        ["order", "POST", "/api/orders", "주문 생성", "Y", "productId, quantity", "Order", "OrderController"],
        ["order", "GET", "/api/orders/{id}", "주문 조회", "Y", "-", "Order", "OrderController"],
        ["order", "GET", "/api/orders", "주문 목록 조회", "Y", "-", "Page<Order>", "OrderController"],
        ["order", "PUT", "/api/orders/{id}/status", "주문 상태 변경", "Y(ADMIN)", "status", "Order", "OrderController"],
        ["order", "POST", "/api/orders/{id}/cancel", "주문 취소", "Y", "-", "Order", "OrderController"],
        # Payment
        ["payment", "POST", "/api/payments/confirm", "결제 승인", "Y", "paymentKey, orderId, amount", "Payment", "PaymentController"],
        ["payment", "POST", "/api/payments/{id}/cancel", "결제 취소", "Y", "cancelReason", "Payment", "PaymentController"],
        ["payment", "POST", "/api/payments/{id}/refund", "환불", "Y", "Idempotency-Key header", "Payment", "PaymentController"],
        # Settlement
        ["settlement", "GET", "/api/settlements", "정산 목록 조회", "Y(ADMIN)", "-", "Page<Settlement>", "SettlementController"],
        ["settlement", "GET", "/api/settlements/{id}", "정산 상세 조회", "Y(ADMIN)", "-", "Settlement", "SettlementController"],
        ["settlement", "GET", "/api/settlements/search", "정산 검색 (ES)", "Y(ADMIN)", "keyword, dateRange, status", "Page<Settlement>", "SettlementSearchController"],
        ["settlement", "GET", "/api/settlements/{id}/pdf", "정산 PDF 다운로드", "Y(ADMIN)", "-", "PDF File", "SettlementController"],
        # Coupon
        ["coupon", "POST", "/api/coupons", "쿠폰 생성", "Y(ADMIN)", "code, type, discountValue, expiresAt", "Coupon", "CouponController"],
        ["coupon", "GET", "/api/coupons", "쿠폰 목록 조회", "Y", "-", "List<Coupon>", "CouponController"],
        ["coupon", "POST", "/api/coupons/{code}/apply", "쿠폰 적용", "Y", "orderId", "Discount", "CouponController"],
        # Review
        ["review", "POST", "/api/reviews", "리뷰 작성", "Y", "productId, rating, content", "Review", "ReviewController"],
        ["review", "GET", "/api/products/{id}/reviews", "상품 리뷰 조회", "N", "-", "List<Review>", "ReviewController"],
    ]
    write_sheet(ws, headers, rows, "API 명세서")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "설계", "API 명세서.xlsx"))
    print("  [OK] API 명세서.xlsx")


def gen_프로그램_목록():
    wb = new_wb()
    ws = wb.active
    headers = ["No", "도메인", "레이어", "프로그램명", "유형", "파일 경로", "설명"]
    programs = [
        # User
        [1, "user", "adapter/in", "AuthController", "Controller", "user/adapter/in/web/AuthController.java", "인증 API (로그인/회원가입)"],
        [2, "user", "adapter/in", "UserController", "Controller", "user/adapter/in/web/UserController.java", "사용자 관리 API"],
        [3, "user", "application", "CreateUserService", "Service", "user/application/service/CreateUserService.java", "회원가입 처리"],
        [4, "user", "application", "LoginService", "Service", "user/application/service/LoginService.java", "로그인 처리"],
        [5, "user", "application", "GetUserService", "Service", "user/application/service/GetUserService.java", "사용자 조회"],
        [6, "user", "application", "PasswordResetService", "Service", "user/application/service/PasswordResetService.java", "비밀번호 재설정"],
        [7, "user", "adapter/out", "UserPersistenceAdapter", "Adapter", "user/adapter/out/persistence/", "사용자 영속성 어댑터"],
        # Product
        [8, "product", "adapter/in", "ProductController", "Controller", "product/adapter/in/web/ProductController.java", "상품 관리 API"],
        [9, "product", "adapter/in", "ProductImageController", "Controller", "product/adapter/in/web/ProductImageController.java", "상품 이미지 API"],
        [10, "product", "adapter/in", "CategoryController", "Controller", "product/adapter/in/web/CategoryController.java", "상품 카테고리 API"],
        [11, "product", "adapter/in", "TagController", "Controller", "product/adapter/in/web/TagController.java", "태그 관리 API"],
        [12, "product", "application", "CreateProductService", "Service", "product/application/service/CreateProductService.java", "상품 등록"],
        [13, "product", "application", "GetProductService", "Service", "product/application/service/GetProductService.java", "상품 조회"],
        [14, "product", "application", "UpdateProductService", "Service", "product/application/service/UpdateProductService.java", "상품 수정"],
        [15, "product", "application", "ManageProductStatusService", "Service", "product/application/service/ManageProductStatusService.java", "상품 상태 관리"],
        [16, "product", "application", "ProductImageService", "Service", "product/application/service/ProductImageService.java", "상품 이미지 처리"],
        [17, "product", "application", "FileStorageService", "Service", "product/application/service/FileStorageService.java", "파일 저장"],
        [18, "product", "application", "CategoryService", "Service", "product/application/service/CategoryService.java", "카테고리 서비스"],
        [19, "product", "application", "TagService", "Service", "product/application/service/TagService.java", "태그 서비스"],
        # Category
        [20, "category", "adapter/in", "AdminEcommerceCategoryController", "Controller", "category/adapter/in/web/AdminEcommerceCategoryController.java", "관리자 카테고리 API"],
        [21, "category", "adapter/in", "PublicEcommerceCategoryController", "Controller", "category/adapter/in/web/PublicEcommerceCategoryController.java", "공개 카테고리 API"],
        [22, "category", "application", "EcommerceCategoryService", "Service", "category/application/service/EcommerceCategoryService.java", "이커머스 카테고리"],
        # Order
        [23, "order", "adapter/in", "OrderController", "Controller", "order/adapter/in/web/OrderController.java", "주문 관리 API"],
        [24, "order", "application", "CreateOrderService", "Service", "order/application/service/CreateOrderService.java", "주문 생성"],
        [25, "order", "application", "ChangeOrderStatusService", "Service", "order/application/service/ChangeOrderStatusService.java", "주문 상태 변경"],
        [26, "order", "application", "GetOrderService", "Service", "order/application/service/GetOrderService.java", "주문 조회"],
        # Payment
        [27, "payment", "adapter/in", "PaymentController", "Controller", "payment/adapter/in/api/PaymentController.java", "결제 관리 API"],
        # Settlement
        [28, "settlement", "adapter/in", "SettlementController", "Controller", "settlement/adapter/in/web/SettlementController.java", "정산 관리 API"],
        [29, "settlement", "adapter/in", "SettlementSearchController", "Controller", "settlement/adapter/in/web/SettlementSearchController.java", "정산 검색 API"],
        [30, "settlement", "application", "CreateDailySettlementsService", "Service", "settlement/application/service/CreateDailySettlementsService.java", "일일 정산 생성"],
        [31, "settlement", "application", "ConfirmDailySettlementsService", "Service", "settlement/application/service/ConfirmDailySettlementsService.java", "정산 확정"],
        [32, "settlement", "application", "AdjustSettlementForRefundService", "Service", "settlement/application/service/AdjustSettlementForRefundService.java", "환불 정산 조정"],
        [33, "settlement", "application", "CreateSettlementFromPaymentService", "Service", "settlement/application/service/CreateSettlementFromPaymentService.java", "결제→정산 생성"],
        [34, "settlement", "application", "GenerateSettlementPdfService", "Service", "settlement/application/service/GenerateSettlementPdfService.java", "정산 PDF 생성"],
        [35, "settlement", "application", "GetSettlementService", "Service", "settlement/application/service/GetSettlementService.java", "정산 조회"],
        [36, "settlement", "application", "IndexSettlementService", "Service", "settlement/application/service/IndexSettlementService.java", "정산 ES 인덱싱"],
        # Coupon
        [37, "coupon", "adapter/in", "CouponController", "Controller", "coupon/adapter/in/web/CouponController.java", "쿠폰 관리 API"],
        [38, "coupon", "application", "CouponService", "Service", "coupon/application/service/CouponService.java", "쿠폰 서비스"],
        # Review
        [39, "review", "adapter/in", "ReviewController", "Controller", "review/adapter/in/web/ReviewController.java", "리뷰 관리 API"],
        # Common
        [40, "common", "config", "RootController", "Controller", "common/config/RootController.java", "루트 경로 핸들러"],
    ]
    write_sheet(ws, headers, programs, "프로그램 목록")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "설계", "프로그램 목록.xlsx"))
    print("  [OK] 프로그램 목록.xlsx")


def gen_공통코드_정의서():
    wb = new_wb()
    ws = wb.active
    headers = ["코드 그룹", "코드 값", "코드명", "설명", "사용 위치"]
    rows = [
        ["OrderStatus", "PENDING", "대기", "주문 생성 초기 상태", "Order.status"],
        ["OrderStatus", "CONFIRMED", "확인", "주문 확인됨", "Order.status"],
        ["OrderStatus", "PAID", "결제완료", "결제 완료된 주문", "Order.status"],
        ["OrderStatus", "SHIPPED", "배송중", "배송 중", "Order.status"],
        ["OrderStatus", "DELIVERED", "배송완료", "배송 완료", "Order.status"],
        ["OrderStatus", "CANCELLED", "취소", "주문 취소됨", "Order.status"],
        ["PaymentStatus", "PENDING", "대기", "결제 대기", "Payment.status"],
        ["PaymentStatus", "AUTHORIZED", "승인", "결제 승인됨", "Payment.status"],
        ["PaymentStatus", "CAPTURED", "매입", "결제 매입 완료", "Payment.status"],
        ["PaymentStatus", "CANCELLED", "취소", "결제 취소됨", "Payment.status"],
        ["PaymentStatus", "REFUNDED", "환불", "환불 완료", "Payment.status"],
        ["PaymentStatus", "FAILED", "실패", "결제 실패", "Payment.status"],
        ["SettlementStatus", "PENDING", "대기", "정산 생성 초기 상태", "Settlement.status"],
        ["SettlementStatus", "CONFIRMED", "확정", "정산 확정됨", "Settlement.status"],
        ["SettlementStatus", "PAID_OUT", "지급완료", "정산금 지급 완료", "Settlement.status"],
        ["SettlementStatus", "ADJUSTED", "조정", "환불 등으로 조정됨", "Settlement.status"],
        ["ProductStatus", "ACTIVE", "활성", "판매 중", "Product.status"],
        ["ProductStatus", "INACTIVE", "비활성", "판매 중지", "Product.status"],
        ["ProductStatus", "OUT_OF_STOCK", "품절", "재고 소진", "Product.status"],
        ["UserRole", "USER", "일반 사용자", "일반 고객 역할", "User.role"],
        ["UserRole", "ADMIN", "관리자", "시스템 관리자 역할", "User.role"],
        ["CouponType", "PERCENTAGE", "정률 할인", "퍼센트 할인", "Coupon.type"],
        ["CouponType", "FIXED_AMOUNT", "정액 할인", "고정 금액 할인", "Coupon.type"],
    ]
    write_sheet(ws, headers, rows, "공통코드 정의서")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "설계", "공통코드 정의서.xlsx"))
    print("  [OK] 공통코드 정의서.xlsx")


def gen_설계단계_완료보고서():
    prs = new_pptx("설계단계 완료보고서", f"{PROJECT_NAME}\n{PROJECT_DATE}")
    add_slide(prs, "1. 설계 개요", [
        "설계 기간: 2026-02-01 ~ 2026-02-28",
        "아키텍처: 헥사고날 (Ports & Adapters)",
        "기술 스택: Spring Boot 3.x + React 18 + PostgreSQL 17 + ES 8.x",
        "총 프로그램 수: 40개 (Controller 16, Service 24)",
    ])
    add_slide(prs, "2. 주요 산출물", [
        "시스템 아키텍처 설계서 — 헥사고날 구조, 인프라, 보안 설계",
        "화면설계서 (UI/UX) — 24개 화면, 네비게이션 구조",
        "데이터베이스 설계서 — 15개 테이블, Flyway V1~V21",
        "API 명세서 — 36개 REST API 엔드포인트",
        "프로그램 목록 — 40개 프로그램 (10개 도메인)",
        "공통코드 정의서 — 23개 코드 값",
    ])
    add_slide(prs, "3. 설계 핵심 결정사항", [
        "동시성 제어: PESSIMISTIC_WRITE 락 + Idempotency-Key",
        "정산 배치: Spring Batch (02:00 생성, 03:00 확정, 03:10 조정)",
        "검색: Elasticsearch + Nori 한글 분석기",
        "도메인↔엔티티 분리: MapStruct Mapper 적용",
        "인증: JWT (Access + Refresh Token)",
    ])
    add_slide(prs, "4. 다음 단계", [
        "구현 단계 진입 (2026-03-01 ~)",
        "도메인별 단위 테스트 작성",
        "통합 테스트 시나리오 실행",
        "API 문서 자동화 (Swagger) 적용",
        "성능 테스트 환경 구축",
    ])
    prs.save(os.path.join(DOCS, "설계", "설계단계 완료보고서.pptx"))
    print("  [OK] 설계단계 완료보고서.pptx")


# ═══════════════════════════════════════════════════════════════
# 3. 구현 단계
# ═══════════════════════════════════════════════════════════════

def gen_단위테스트_결과서():
    wb = new_wb()
    ws = wb.active
    headers = ["No", "테스트 파일", "테스트 메서드", "테스트 대상", "결과", "비고"]
    rows = [
        [1, "CreateUserServiceTest", "회원가입_성공", "CreateUserService", "PASS", ""],
        [2, "CreateUserServiceTest", "이메일_중복_실패", "CreateUserService", "PASS", ""],
        [3, "LoginServiceTest", "로그인_성공", "LoginService", "PASS", ""],
        [4, "LoginServiceTest", "잘못된_비밀번호_실패", "LoginService", "PASS", ""],
        [5, "UserTest", "사용자_생성", "User (Domain)", "PASS", ""],
        [6, "UserRoleTest", "역할_검증", "UserRole (Domain)", "PASS", ""],
        [7, "UserPersistenceAdapterTest", "사용자_저장_조회", "UserPersistenceAdapter", "PASS", ""],
        [8, "CreateOrderServiceTest", "주문_생성_성공", "CreateOrderService", "PASS", ""],
        [9, "OrderTest", "주문_도메인_검증", "Order (Domain)", "PASS", ""],
        [10, "ProductTest", "상품_도메인_검증", "Product (Domain)", "PASS", ""],
        [11, "CouponServiceTest", "쿠폰_발급_성공", "CouponService", "PASS", ""],
        [12, "CouponServiceTest", "쿠폰_적용_성공", "CouponService", "PASS", ""],
        [13, "EcommerceCategoryTest", "카테고리_생성", "EcommerceCategory (Domain)", "PASS", ""],
        [14, "CreateDailySettlementsServiceTest", "일일정산_생성", "CreateDailySettlementsService", "PASS", ""],
        [15, "ConfirmDailySettlementsServiceTest", "정산_확정", "ConfirmDailySettlementsService", "PASS", ""],
        [16, "AdjustSettlementForRefundServiceTest", "환불_정산조정", "AdjustSettlementForRefundService", "PASS", ""],
        [17, "SettlementTest", "정산_도메인_검증", "Settlement (Domain)", "PASS", ""],
        [18, "SettlementStatusTest", "정산상태_전이", "SettlementStatus (Domain)", "PASS", ""],
        [19, "SettlementSearchControllerTest", "정산_검색_API", "SettlementSearchController", "PASS", ""],
        [20, "ShoppingFlowIntegrationTest", "주문→결제→정산_E2E", "통합 테스트", "PASS", "H2 인메모리 DB"],
    ]
    write_sheet(ws, headers, rows, "단위테스트 결과서")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "구현", "단위테스트 결과서.xlsx"))
    print("  [OK] 단위테스트 결과서.xlsx")


def gen_코드리뷰_결과서():
    doc = new_doc("코드 리뷰 결과서", "코드 리뷰 수행 기록 및 조치 결과")

    add_heading(doc, "1. 코드 리뷰 개요")
    add_table(doc,
        ["항목", "내용"],
        [
            ["프로젝트", PROJECT_NAME],
            ["리뷰 기간", "2026-03-01 ~ 2026-03-25"],
            ["리뷰 방법", "GitHub PR 기반 코드 리뷰"],
            ["리뷰 대상", "전체 도메인 (user, product, order, payment, settlement, coupon, review, category)"],
        ]
    )

    add_heading(doc, "2. 리뷰 결과 요약")
    add_table(doc,
        ["구분", "건수", "조치 완료", "미조치"],
        [
            ["아키텍처 위반", "3", "3", "0"],
            ["코딩 컨벤션", "5", "5", "0"],
            ["보안 이슈", "2", "2", "0"],
            ["성능 이슈", "1", "1", "0"],
            ["기타", "4", "4", "0"],
            ["합계", "15", "15", "0"],
        ]
    )

    add_heading(doc, "3. 주요 리뷰 항목")
    add_table(doc,
        ["No", "PR #", "도메인", "이슈 유형", "내용", "조치 결과"],
        [
            ["1", "#42", "common", "아키텍처", "도커 컴포즈/프론트/백 연결 구조 개선", "조치 완료"],
            ["2", "#43", "infra", "아키텍처", "k8s 구조 변경", "조치 완료"],
            ["3", "#45", "payment", "보안", "환불 멱등성 키 검증 강화", "Idempotency-Key 적용"],
            ["4", "#46", "settlement", "성능", "정산 배치 쿼리 최적화", "인덱스 추가"],
            ["5", "#47", "product", "컨벤션", "MapStruct Mapper 패턴 통일", "조치 완료"],
            ["6", "#48", "settlement", "아키텍처", "정산 금액 분할 (V9)", "settlement/fee/net_amount 분리"],
            ["7", "#49", "전체", "컨벤션", "도메인 레이어 Spring 의존성 제거", "조치 완료"],
        ]
    )

    add_heading(doc, "4. 리뷰 기준")
    criteria = [
        "헥사고날 아키텍처 의존성 규칙 준수 (domain ← application ← adapter)",
        "domain 레이어에 Spring/JPA 의존성 없음",
        "UseCase/Port 인터페이스를 통한 계층 간 통신",
        "JPA Entity와 Domain 객체 분리 (MapStruct 매핑)",
        "Request/Response DTO는 adapter/in/web/dto/ 위치",
        "OWASP Top 10 보안 취약점 점검",
    ]
    for c in criteria:
        doc.add_paragraph(c, style='List Bullet')

    doc.save(os.path.join(DOCS, "구현", "코드 리뷰 결과서.docx"))
    print("  [OK] 코드 리뷰 결과서.docx")


def gen_코드인스펙션_보고서():
    wb = new_wb()
    ws = wb.active
    headers = ["No", "점검 항목", "도구", "결과", "심각도", "내용", "조치 상태"]
    rows = [
        [1, "정적 분석", "SonarCloud", "PASS", "-", "코드 스멜, 버그, 취약점 0건", "완료"],
        [2, "의존성 취약점", "Snyk", "PASS", "-", "High/Critical 취약점 0건", "완료"],
        [3, "코드 커버리지", "JaCoCo", "WARN", "Medium", "현재 커버리지 약 5.8% (목표 70%)", "진행중"],
        [4, "코딩 표준 (FE)", "ESLint", "PASS", "-", "린트 에러 0건", "완료"],
        [5, "타입 안전성 (FE)", "TypeScript", "PASS", "-", "타입 에러 0건", "완료"],
        [6, "헥사고날 아키텍처", "수동 점검", "PASS", "-", "domain→외부 의존성 없음 확인", "완료"],
        [7, "SQL Injection", "수동 점검", "PASS", "-", "JPA 파라미터 바인딩 사용 확인", "완료"],
        [8, "XSS 방지", "수동 점검", "PASS", "-", "React 자동 이스케이핑 적용", "완료"],
        [9, "인증/인가", "수동 점검", "PASS", "-", "JWT + Spring Security 적용 확인", "완료"],
        [10, "비밀번호 해싱", "수동 점검", "PASS", "-", "BCrypt 적용 확인", "완료"],
    ]
    write_sheet(ws, headers, rows, "코드 인스펙션 보고서")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "구현", "코드 인스펙션 보고서.xlsx"))
    print("  [OK] 코드 인스펙션 보고서.xlsx")


def gen_개발진척_관리표():
    wb = new_wb()
    ws = wb.active
    headers = ["도메인", "기능", "설계", "구현", "단위테스트", "통합테스트", "진행률", "비고"]
    rows = [
        ["user", "회원가입/로그인", "완료", "완료", "완료", "완료", "100%", ""],
        ["user", "비밀번호 재설정", "완료", "완료", "미작성", "-", "70%", "테스트 필요"],
        ["product", "상품 CRUD", "완료", "완료", "미작성", "-", "70%", "테스트 필요"],
        ["product", "상품 검색 (ES)", "완료", "완료", "미작성", "-", "70%", "테스트 필요"],
        ["product", "이미지 관리", "완료", "완료", "미작성", "-", "70%", "테스트 필요"],
        ["category", "카테고리 CRUD", "완료", "완료", "미작성", "-", "70%", "도메인 테스트만 존재"],
        ["order", "주문 생성", "완료", "완료", "완료", "완료", "100%", ""],
        ["order", "주문 상태 변경/취소", "완료", "완료", "미작성", "-", "70%", "테스트 필요"],
        ["payment", "결제 승인/취소", "완료", "완료", "미작성", "-", "70%", "테스트 필요"],
        ["payment", "환불", "완료", "완료", "미작성", "-", "70%", "멱등성 보장 구현"],
        ["settlement", "일일 정산 생성", "완료", "완료", "완료", "완료", "100%", ""],
        ["settlement", "정산 확정", "완료", "완료", "완료", "-", "90%", ""],
        ["settlement", "환불 정산 조정", "완료", "완료", "완료", "-", "90%", ""],
        ["settlement", "정산 검색 (ES)", "완료", "완료", "완료", "-", "90%", ""],
        ["settlement", "정산 PDF", "완료", "완료", "미작성", "-", "70%", "테스트 필요"],
        ["coupon", "쿠폰 발급/적용", "완료", "완료", "완료", "-", "90%", ""],
        ["review", "리뷰 작성/조회", "완료", "완료", "미작성", "-", "70%", "테스트 필요"],
        ["game", "오목/바둑", "미완성", "미완성", "미작성", "-", "20%", "Controller만 존재"],
    ]
    write_sheet(ws, headers, rows, "개발 진척 관리표")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "구현", "개발 진척 관리표.xlsx"))
    print("  [OK] 개발 진척 관리표.xlsx")


def gen_데이터이행_계획서():
    doc = new_doc("데이터 이행 계획서", "기존 데이터 마이그레이션 계획")

    add_heading(doc, "1. 이행 개요")
    doc.add_paragraph("본 문서는 Lemuel 시스템의 데이터 이행 계획을 기술한다. "
                       "Flyway 마이그레이션을 통한 스키마 이행과 시드 데이터 적재를 포함한다.")

    add_heading(doc, "2. 이행 범위")
    add_table(doc,
        ["대상", "데이터 유형", "건수(예상)", "이행 방법"],
        [
            ["스키마", "DDL (V1~V16)", "15개 테이블", "Flyway 자동 마이그레이션"],
            ["시드 데이터", "DML (V17, V18, V21)", "-", "Flyway 시드 마이그레이션"],
            ["사용자 데이터", "기존 사용자 정보", "-", "해당없음 (신규 시스템)"],
            ["상품 데이터", "상품 카탈로그", "-", "시드 데이터 포함"],
        ]
    )

    add_heading(doc, "3. 이행 순서")
    steps = [
        "1단계: PostgreSQL 17 설치 및 opslab 스키마 생성",
        "2단계: Flyway V1~V16 실행 (스키마 생성)",
        "3단계: V17 시드 데이터 적재",
        "4단계: V18 시드 매니저 설정",
        "5단계: V19~V20 추가 테이블 생성 (reviews, coupons)",
        "6단계: V21 2026년 1월 시드 데이터 적재",
        "7단계: Elasticsearch 인덱스 생성 및 데이터 동기화",
    ]
    for s in steps:
        doc.add_paragraph(s, style='List Number')

    add_heading(doc, "4. 이행 환경")
    add_table(doc,
        ["환경", "DB", "ES", "비고"],
        [
            ["개발", "PostgreSQL 17 (Docker)", "ES 8.x (Docker)", "docker compose up"],
            ["스테이징", "PostgreSQL 17", "ES 8.x", "CI/CD 자동 배포"],
            ["운영", "PostgreSQL 17", "ES 8.x", "수동 승인 후 배포"],
        ]
    )

    add_heading(doc, "5. 롤백 계획")
    doc.add_paragraph("Flyway는 유료 버전에서만 롤백을 지원하므로, 각 마이그레이션 전 DB 백업을 수행한다.")
    rollback_steps = [
        "마이그레이션 실행 전 pg_dump로 전체 백업",
        "문제 발생 시 pg_restore로 복원",
        "Flyway flyway_schema_history 테이블에서 실패 레코드 삭제",
    ]
    for s in rollback_steps:
        doc.add_paragraph(s, style='List Bullet')

    doc.save(os.path.join(DOCS, "구현", "데이터 이행 계획서.docx"))
    print("  [OK] 데이터 이행 계획서.docx")


def gen_데이터이행_결과서():
    doc = new_doc("데이터 이행 결과서", "데이터 이행 수행 결과 및 검증")

    add_heading(doc, "1. 이행 수행 결과")
    add_table(doc,
        ["단계", "작업", "수행일", "결과", "비고"],
        [
            ["1", "PostgreSQL 설치/스키마 생성", "2026-03-01", "성공", "opslab 스키마"],
            ["2", "Flyway V1~V16 실행", "2026-03-01", "성공", "15개 테이블 생성"],
            ["3", "V17 시드 데이터", "2026-03-01", "성공", "기본 데이터 적재"],
            ["4", "V18 시드 매니저", "2026-03-01", "성공", ""],
            ["5", "V19~V20 테이블 추가", "2026-03-10", "성공", "reviews, coupons"],
            ["6", "V21 시드 데이터", "2026-03-15", "성공", "2026년 1월 데이터"],
            ["7", "ES 인덱스 동기화", "2026-03-15", "성공", "상품/정산 인덱스"],
        ]
    )

    add_heading(doc, "2. 데이터 검증 결과")
    add_table(doc,
        ["검증 항목", "방법", "결과", "비고"],
        [
            ["테이블 수", "information_schema 조회", "15개 — 정상", ""],
            ["마이그레이션 이력", "flyway_schema_history 조회", "V1~V21 모두 성공", ""],
            ["인덱스", "pg_indexes 조회", "V3 인덱스 모두 생성됨", ""],
            ["FK 제약조건", "information_schema 조회", "모든 FK 정상", ""],
            ["시드 데이터", "SELECT COUNT 검증", "정상 적재", ""],
            ["ES 인덱스", "GET _cat/indices", "settlements, products 인덱스 정상", ""],
        ]
    )

    add_heading(doc, "3. 이슈 및 조치")
    doc.add_paragraph("V16: PG 트랜잭션 ID 길이 부족 → VARCHAR(200)으로 확장하여 해결")

    doc.save(os.path.join(DOCS, "구현", "데이터 이행 결과서.docx"))
    print("  [OK] 데이터 이행 결과서.docx")


def gen_구현단계_완료보고서():
    prs = new_pptx("구현단계 완료보고서", f"{PROJECT_NAME}\n{PROJECT_DATE}")
    add_slide(prs, "1. 구현 개요", [
        "구현 기간: 2026-03-01 ~ 2026-03-25",
        "총 소스 파일: 294개 (백엔드), 프론트엔드 40+개",
        "총 테스트: 17개 파일, 20+ 테스트 케이스",
        "Flyway 마이그레이션: V1~V21 (21개)",
    ])
    add_slide(prs, "2. 도메인별 구현 현황", [
        "user: 회원가입, 로그인, 비밀번호 재설정 — 100% 구현",
        "product: 상품 CRUD, 검색, 이미지 — 100% 구현",
        "order: 주문 생성, 상태변경, 취소 — 100% 구현",
        "payment: Toss PG 연동, 환불 (멱등성) — 100% 구현",
        "settlement: 배치 정산, ES 검색, PDF — 100% 구현",
        "coupon: 쿠폰 발급/적용 — 100% 구현",
        "game: 오목/바둑 — 20% (Controller만)",
    ])
    add_slide(prs, "3. 코드 품질", [
        "정적 분석 (SonarCloud): PASS",
        "보안 스캔 (Snyk): High/Critical 0건",
        "ESLint / TypeScript: 에러 0건",
        "테스트 커버리지: 약 5.8% (목표 70% 미달 — 보완 필요)",
    ])
    add_slide(prs, "4. 이슈 및 리스크", [
        "[Critical] 테스트 커버리지 심각하게 부족 (17/294 파일)",
        "[Medium] API 문서화 미비 (Swagger 미적용)",
        "[Low] 게임 도메인 미완성",
        "[Action] 서비스/컨트롤러 단위 테스트 대량 작성 필요",
    ])
    add_slide(prs, "5. 다음 단계", [
        "테스트 단계 진입",
        "통합 테스트 시나리오 확대",
        "성능 테스트 수행",
        "보안 점검 강화",
        "UAT 수행",
    ])
    prs.save(os.path.join(DOCS, "구현", "구현단계 완료보고서.pptx"))
    print("  [OK] 구현단계 완료보고서.pptx")


# ═══════════════════════════════════════════════════════════════
# 4. 테스트 단계
# ═══════════════════════════════════════════════════════════════

def gen_통합테스트_계획서():
    doc = new_doc("통합테스트 계획서", "통합테스트 전략/범위/일정/환경 정의")

    add_heading(doc, "1. 테스트 개요")
    add_table(doc,
        ["항목", "내용"],
        [
            ["프로젝트", PROJECT_NAME],
            ["테스트 기간", "2026-03-25 ~ 2026-04-10"],
            ["테스트 환경", "H2 인메모리 DB (단위), PostgreSQL 17 (통합)"],
            ["테스트 도구", "JUnit 5, Vitest, MSW, Testing Library"],
            ["테스트 전략", "기능간 연동 중심 테스트"],
        ]
    )

    add_heading(doc, "2. 테스트 범위")
    add_table(doc,
        ["No", "테스트 영역", "연동 대상", "우선순위"],
        [
            [1, "주문 → 결제 연동", "Order ↔ Payment", "필수"],
            [2, "결제 → 정산 연동", "Payment ↔ Settlement", "필수"],
            [3, "주문 → 결제 → 정산 E2E", "Order → Payment → Settlement", "필수"],
            [4, "환불 → 정산조정 연동", "Payment(Refund) → Settlement(Adjust)", "필수"],
            [5, "상품 → 주문 재고 연동", "Product(Stock) ↔ Order", "필수"],
            [6, "쿠폰 → 주문 할인 연동", "Coupon ↔ Order", "높음"],
            [7, "사용자 → 주문 인증 연동", "User(JWT) ↔ Order", "높음"],
            [8, "상품 → ES 검색 연동", "Product ↔ Elasticsearch", "보통"],
            [9, "정산 → ES 검색 연동", "Settlement ↔ Elasticsearch", "보통"],
            [10, "정산 → PDF 생성", "Settlement → iText PDF", "보통"],
        ]
    )

    add_heading(doc, "3. 테스트 환경")
    add_table(doc,
        ["환경", "DB", "ES", "PG", "비고"],
        [
            ["단위 테스트", "H2 인메모리", "Mock", "Mock", "JUnit 5"],
            ["통합 테스트", "PostgreSQL 17", "ES 8.x", "Toss Sandbox", "Docker Compose"],
            ["E2E 테스트", "PostgreSQL 17", "ES 8.x", "Toss Sandbox", "프론트 + 백엔드"],
        ]
    )

    add_heading(doc, "4. 합격 기준")
    criteria = [
        "모든 필수 테스트 시나리오 PASS",
        "결함 심각도 Critical/High: 0건",
        "테스트 커버리지 70% 이상 (JaCoCo)",
        "응답 시간 95th percentile 500ms 이내",
    ]
    for c in criteria:
        doc.add_paragraph(c, style='List Bullet')

    doc.save(os.path.join(DOCS, "테스트", "통합테스트 계획서.docx"))
    print("  [OK] 통합테스트 계획서.docx")


def gen_통합테스트_시나리오():
    wb = new_wb()
    ws = wb.active
    headers = ["시나리오 ID", "시나리오명", "사전조건", "테스트 단계", "예상 결과", "우선순위"]
    rows = [
        ["IT-001", "주문→결제→정산 E2E", "상품 존재, 사용자 로그인", "1.주문생성 2.Toss결제 3.일일정산배치 4.정산확정", "CONFIRMED 정산 생성", "필수"],
        ["IT-002", "환불→정산조정", "결제 CAPTURED 상태", "1.환불요청 2.PG취소 3.정산조정생성", "ADJUSTED 정산 조정 생성", "필수"],
        ["IT-003", "환불 멱등성 검증", "결제 CAPTURED 상태", "1.환불요청(Key=A) 2.동일요청(Key=A)", "두 번째 요청은 기존 결과 반환", "필수"],
        ["IT-004", "환불 동시성 검증", "결제 CAPTURED 상태", "1.동시 환불 요청 2건", "1건만 성공, 1건 실패(락)", "필수"],
        ["IT-005", "주문-재고 연동", "상품 재고 10개", "1.5개 주문 2.재고확인 3.15개 주문(실패)", "재고 5개, InsufficientStockException", "필수"],
        ["IT-006", "쿠폰 적용 주문", "유효한 쿠폰 존재", "1.주문생성 2.쿠폰적용 3.결제금액확인", "할인 적용된 결제 금액", "높음"],
        ["IT-007", "만료 쿠폰 사용 시도", "만료된 쿠폰", "1.주문생성 2.만료쿠폰적용", "오류 반환", "높음"],
        ["IT-008", "상품 ES 검색", "상품 데이터 인덱싱 완료", "1.한글 키워드 검색 2.결과 확인", "Nori 분석기로 검색 결과 반환", "보통"],
        ["IT-009", "정산 ES 검색", "정산 데이터 인덱싱 완료", "1.기간/상태 검색 2.결과 확인", "필터 조건에 맞는 정산 목록", "보통"],
        ["IT-010", "정산 PDF 다운로드", "CONFIRMED 정산 존재", "1.PDF 생성 요청 2.파일 다운로드", "유효한 PDF 파일", "보통"],
        ["IT-011", "배치 중복 실행 방지", "오늘 배치 이미 실행됨", "1.동일 파라미터로 재실행", "Job 스킵됨", "필수"],
        ["IT-012", "JWT 인증 검증", "유효한 JWT 토큰", "1.인증 필요 API 호출 2.만료 토큰으로 호출", "200 OK / 401 Unauthorized", "높음"],
    ]
    write_sheet(ws, headers, rows, "통합테스트 시나리오")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "테스트", "통합테스트 시나리오.xlsx"))
    print("  [OK] 통합테스트 시나리오.xlsx")


def gen_통합테스트_결과서():
    wb = new_wb()
    ws = wb.active
    headers = ["시나리오 ID", "시나리오명", "수행일", "결과", "결함 ID", "비고"]
    rows = [
        ["IT-001", "주문→결제→정산 E2E", "2026-03-25", "PASS", "-", "ShoppingFlowIntegrationTest"],
        ["IT-002", "환불→정산조정", "2026-03-25", "PASS", "-", "AdjustSettlementForRefundServiceTest"],
        ["IT-003", "환불 멱등성 검증", "-", "미수행", "-", "테스트 작성 필요"],
        ["IT-004", "환불 동시성 검증", "-", "미수행", "-", "테스트 작성 필요"],
        ["IT-005", "주문-재고 연동", "-", "미수행", "-", "테스트 작성 필요"],
        ["IT-006", "쿠폰 적용 주문", "-", "미수행", "-", "테스트 작성 필요"],
        ["IT-007", "만료 쿠폰 사용 시도", "-", "미수행", "-", "테스트 작성 필요"],
        ["IT-008", "상품 ES 검색", "-", "미수행", "-", "ES 환경 필요"],
        ["IT-009", "정산 ES 검색", "2026-03-25", "PASS", "-", "SettlementSearchControllerTest"],
        ["IT-010", "정산 PDF 다운로드", "-", "미수행", "-", "테스트 작성 필요"],
        ["IT-011", "배치 중복 실행 방지", "-", "미수행", "-", "테스트 작성 필요"],
        ["IT-012", "JWT 인증 검증", "-", "미수행", "-", "테스트 작성 필요"],
    ]
    write_sheet(ws, headers, rows, "통합테스트 결과서")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "테스트", "통합테스트 결과서.xlsx"))
    print("  [OK] 통합테스트 결과서.xlsx")


def gen_성능테스트_계획서():
    doc = new_doc("성능테스트 계획서", "성능 목표/시나리오/도구 정의")

    add_heading(doc, "1. 테스트 개요")
    add_table(doc,
        ["항목", "내용"],
        [
            ["프로젝트", PROJECT_NAME],
            ["테스트 기간", "2026-04-01 ~ 2026-04-05"],
            ["테스트 도구", "JMeter / k6"],
            ["테스트 환경", "Docker Compose (PostgreSQL 17, ES 8.x)"],
        ]
    )

    add_heading(doc, "2. 성능 목표")
    add_table(doc,
        ["지표", "목표값", "비고"],
        [
            ["응답 시간 (평균)", "200ms 이내", ""],
            ["응답 시간 (P95)", "500ms 이내", "NFR-PERF-001"],
            ["응답 시간 (P99)", "1000ms 이내", ""],
            ["동시 사용자", "500명 이상", "NFR-PERF-002"],
            ["TPS", "100 TPS 이상", "주요 API 기준"],
            ["에러율", "0.1% 미만", ""],
            ["CPU 사용률", "80% 미만", "피크 시"],
            ["메모리 사용률", "80% 미만", "피크 시"],
        ]
    )

    add_heading(doc, "3. 테스트 시나리오")
    add_table(doc,
        ["No", "시나리오", "동시 사용자", "Duration", "대상 API"],
        [
            [1, "상품 목록 조회 부하", "100 → 500", "10분", "GET /api/products"],
            [2, "상품 검색 부하", "50 → 200", "10분", "GET /api/products/search"],
            [3, "주문 생성 부하", "50 → 200", "10분", "POST /api/orders"],
            [4, "결제 승인 부하", "50 → 100", "10분", "POST /api/payments/confirm"],
            [5, "정산 조회 부하", "50 → 200", "10분", "GET /api/settlements"],
            [6, "정산 검색 (ES) 부하", "50 → 200", "10분", "GET /api/settlements/search"],
            [7, "종합 시나리오", "100 → 500", "30분", "주문→결제 흐름"],
        ]
    )

    doc.save(os.path.join(DOCS, "테스트", "성능테스트 계획서.docx"))
    print("  [OK] 성능테스트 계획서.docx")


def gen_성능테스트_결과서():
    doc = new_doc("성능테스트 결과서", "TPS/응답시간/자원사용률 등 결과")

    add_heading(doc, "1. 테스트 수행 요약")
    doc.add_paragraph("※ 본 문서는 성능테스트 수행 후 실측 데이터로 갱신할 템플릿입니다.")

    add_heading(doc, "2. 테스트 결과")
    add_table(doc,
        ["시나리오", "TPS", "평균 응답(ms)", "P95(ms)", "P99(ms)", "에러율", "결과"],
        [
            ["상품 목록 조회", "-", "-", "-", "-", "-", "미수행"],
            ["상품 검색 (ES)", "-", "-", "-", "-", "-", "미수행"],
            ["주문 생성", "-", "-", "-", "-", "-", "미수행"],
            ["결제 승인", "-", "-", "-", "-", "-", "미수행"],
            ["정산 조회", "-", "-", "-", "-", "-", "미수행"],
            ["정산 검색 (ES)", "-", "-", "-", "-", "-", "미수행"],
            ["종합 시나리오", "-", "-", "-", "-", "-", "미수행"],
        ]
    )

    add_heading(doc, "3. 자원 사용률")
    add_table(doc,
        ["리소스", "평균", "최대", "목표", "결과"],
        [
            ["CPU (Backend)", "-", "-", "< 80%", "미수행"],
            ["Memory (Backend)", "-", "-", "< 80%", "미수행"],
            ["CPU (DB)", "-", "-", "< 80%", "미수행"],
            ["DB Connection Pool", "-", "-", "< 80%", "미수행"],
        ]
    )

    add_heading(doc, "4. 개선 권고사항")
    doc.add_paragraph("성능테스트 수행 후 작성 예정")

    doc.save(os.path.join(DOCS, "테스트", "성능테스트 결과서.docx"))
    print("  [OK] 성능테스트 결과서.docx")


def gen_보안점검_결과서():
    doc = new_doc("보안 점검 결과서", "보안 취약점 점검 및 조치 결과")

    add_heading(doc, "1. 점검 개요")
    add_table(doc,
        ["항목", "내용"],
        [
            ["점검 기간", "2026-03-20 ~ 2026-03-25"],
            ["점검 도구", "Snyk, SonarCloud, 수동 점검"],
            ["점검 기준", "OWASP Top 10 (2021)"],
        ]
    )

    add_heading(doc, "2. OWASP Top 10 점검 결과")
    add_table(doc,
        ["No", "항목", "점검 내용", "결과", "조치"],
        [
            ["A01", "Broken Access Control", "API 인가 검증", "양호", "Spring Security + JWT 역할 기반 인가"],
            ["A02", "Cryptographic Failures", "비밀번호/토큰 암호화", "양호", "BCrypt 해싱, JWT 서명"],
            ["A03", "Injection", "SQL Injection, XSS", "양호", "JPA 파라미터 바인딩, React 자동 이스케이핑"],
            ["A04", "Insecure Design", "아키텍처 보안", "양호", "헥사고날 아키텍처, 도메인 분리"],
            ["A05", "Security Misconfiguration", "서버/DB 설정", "주의", "환경변수 관리 필요 (.env)"],
            ["A06", "Vulnerable Components", "의존성 취약점", "양호", "Snyk 스캔 통과"],
            ["A07", "Auth Failures", "인증 메커니즘", "양호", "JWT + Refresh Token"],
            ["A08", "Data Integrity Failures", "데이터 무결성", "양호", "FK 제약조건, 트랜잭션 관리"],
            ["A09", "Logging Failures", "로깅/모니터링", "주의", "모니터링 체계 구축 필요"],
            ["A10", "SSRF", "서버측 요청 위조", "양호", "외부 URL 직접 호출 없음"],
        ]
    )

    add_heading(doc, "3. 결제/정산 보안 점검")
    add_table(doc,
        ["점검 항목", "결과", "상세"],
        [
            ["환불 멱등성", "양호", "Idempotency-Key 헤더 기반 중복 방지"],
            ["환불 동시성", "양호", "PESSIMISTIC_WRITE 락 적용"],
            ["PG 통신 보안", "양호", "HTTPS + Secret Key 인증"],
            ["정산 데이터 접근 제어", "양호", "ADMIN 역할만 접근 가능"],
            ["배치 중복 실행", "양호", "Spring Batch JobRepository"],
        ]
    )

    add_heading(doc, "4. Snyk 의존성 스캔 결과")
    add_table(doc,
        ["심각도", "건수", "조치"],
        [
            ["Critical", "0", "-"],
            ["High", "0", "-"],
            ["Medium", "0", "-"],
            ["Low", "0", "-"],
        ]
    )

    add_heading(doc, "5. 개선 권고사항")
    recommendations = [
        "[주의] 환경변수(.env) 관리: Secret 값을 Vault 등으로 이관 권장",
        "[주의] 로깅/모니터링: APM 도구(Grafana, Prometheus) 도입 권장",
        "[권장] API Rate Limiting 적용",
        "[권장] CORS 정책 강화",
    ]
    for r in recommendations:
        doc.add_paragraph(r, style='List Bullet')

    doc.save(os.path.join(DOCS, "테스트", "보안 점검 결과서.docx"))
    print("  [OK] 보안 점검 결과서.docx")


def gen_UAT_계획서():
    doc = new_doc("사용자 인수테스트(UAT) 계획서", "최종사용자 인수 테스트 계획")

    add_heading(doc, "1. UAT 개요")
    add_table(doc,
        ["항목", "내용"],
        [
            ["프로젝트", PROJECT_NAME],
            ["테스트 기간", "2026-04-10 ~ 2026-04-17"],
            ["테스트 환경", "스테이징 환경 (Docker Compose)"],
            ["참여자", "기획자, PM, 최종 사용자 대표"],
        ]
    )

    add_heading(doc, "2. UAT 범위")
    add_table(doc,
        ["No", "테스트 영역", "시나리오 수", "테스터"],
        [
            [1, "회원가입/로그인", 3, "사용자 대표"],
            [2, "상품 검색/조회", 4, "사용자 대표"],
            [3, "장바구니/주문", 5, "사용자 대표"],
            [4, "결제 (Toss)", 3, "사용자 대표"],
            [5, "마이페이지/주문내역", 3, "사용자 대표"],
            [6, "환불", 2, "사용자 대표"],
            [7, "관리자-상품관리", 4, "기획자"],
            [8, "관리자-정산관리", 5, "기획자"],
            [9, "관리자-카테고리", 3, "기획자"],
            [10, "관리자-쿠폰관리", 3, "기획자"],
        ]
    )

    add_heading(doc, "3. 합격 기준")
    criteria = [
        "전체 시나리오 90% 이상 PASS",
        "심각도 Critical 결함: 0건",
        "심각도 High 결함: 해결 완료",
        "주요 업무 흐름(주문→결제→정산) 정상 동작",
    ]
    for c in criteria:
        doc.add_paragraph(c, style='List Bullet')

    doc.save(os.path.join(DOCS, "테스트", "사용자 인수테스트(UAT) 계획서.docx"))
    print("  [OK] 사용자 인수테스트(UAT) 계획서.docx")


def gen_UAT_결과서():
    wb = new_wb()
    ws = wb.active
    headers = ["No", "테스트 영역", "시나리오", "수행일", "결과", "결함 ID", "비고"]
    rows = [
        [1, "회원가입", "이메일/비밀번호로 회원가입", "-", "미수행", "-", ""],
        [2, "로그인", "로그인 후 JWT 토큰 발급 확인", "-", "미수행", "-", ""],
        [3, "비밀번호 재설정", "이메일 발송 → 비밀번호 변경", "-", "미수행", "-", ""],
        [4, "상품 검색", "키워드 검색, 카테고리 필터", "-", "미수행", "-", ""],
        [5, "상품 상세", "상품 정보, 이미지, 리뷰 확인", "-", "미수행", "-", ""],
        [6, "장바구니", "상품 담기/수량변경/삭제", "-", "미수행", "-", ""],
        [7, "쿠폰 적용", "유효한 쿠폰 코드 입력 → 할인 확인", "-", "미수행", "-", ""],
        [8, "주문 생성", "주문 정보 확인 → 결제 진행", "-", "미수행", "-", ""],
        [9, "Toss 결제", "결제 페이지 → 결제 완료", "-", "미수행", "-", "Sandbox"],
        [10, "주문 내역", "마이페이지에서 주문 목록 확인", "-", "미수행", "-", ""],
        [11, "환불 요청", "결제 완료 주문 환불", "-", "미수행", "-", ""],
        [12, "관리자 로그인", "관리자 계정 로그인", "-", "미수행", "-", ""],
        [13, "상품 관리", "상품 등록/수정/삭제", "-", "미수행", "-", ""],
        [14, "카테고리 관리", "카테고리 추가/수정/삭제", "-", "미수행", "-", ""],
        [15, "정산 대시보드", "정산 현황 조회, 기간 검색", "-", "미수행", "-", ""],
        [16, "정산 PDF", "정산 내역 PDF 다운로드", "-", "미수행", "-", ""],
    ]
    write_sheet(ws, headers, rows, "UAT 결과서")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "테스트", "UAT 결과서.xlsx"))
    print("  [OK] UAT 결과서.xlsx")


def gen_결함관리_대장():
    wb = new_wb()
    ws = wb.active
    headers = ["결함 ID", "발견일", "심각도", "도메인", "제목", "설명", "발견 단계", "상태", "조치일", "조치 내용"]
    rows = [
        ["DEF-001", "2026-03-10", "Medium", "payment", "PG 트랜잭션 ID 길이 부족", "VARCHAR(100)으로 Toss 긴 트랜잭션 ID 저장 불가", "구현", "종결", "2026-03-10", "V16 마이그레이션으로 VARCHAR(200) 확장"],
        ["DEF-002", "2026-03-15", "Low", "settlement", "정산 금액 단일 컬럼", "수수료/순수익 분리 필요", "코드리뷰", "종결", "2026-03-15", "V9 마이그레이션으로 settlement/fee/net_amount 분리"],
        ["DEF-003", "2026-03-20", "High", "전체", "테스트 커버리지 부족", "17/294 파일 (5.8%), 목표 70% 미달", "코드인스펙션", "미해결", "-", "서비스/컨트롤러 테스트 대량 작성 예정"],
        ["DEF-004", "2026-03-22", "Low", "game", "게임 도메인 미완성", "Controller만 존재, 로직 미구현", "코드리뷰", "보류", "-", "범위 조정 검토"],
        ["DEF-005", "2026-03-25", "Medium", "common", "API 문서 미비", "Swagger/OpenAPI 미적용", "코드인스펙션", "미해결", "-", "Swagger 적용 예정"],
    ]
    write_sheet(ws, headers, rows, "결함 관리 대장")
    auto_width(ws, len(headers))
    wb.save(os.path.join(DOCS, "테스트", "결함 관리 대장.xlsx"))
    print("  [OK] 결함 관리 대장.xlsx")


# ═══════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════

if __name__ == "__main__":
    make_dirs()
    print("\n=== 분석 단계 ===")
    gen_현행시스템_분석서()
    gen_요구사항_정의서()
    gen_RTM()
    gen_업무프로세스_정의서()
    gen_유스케이스_명세서()
    gen_인터뷰_회의록()
    gen_분석단계_완료보고서()

    print("\n=== 설계 단계 ===")
    gen_시스템_아키텍처_설계서()
    gen_화면설계서()
    gen_데이터베이스_설계서()
    gen_테이블_정의서()
    gen_인터페이스_설계서()
    gen_API_명세서()
    gen_프로그램_목록()
    gen_공통코드_정의서()
    gen_설계단계_완료보고서()

    print("\n=== 구현 단계 ===")
    gen_단위테스트_결과서()
    gen_코드리뷰_결과서()
    gen_코드인스펙션_보고서()
    gen_개발진척_관리표()
    gen_데이터이행_계획서()
    gen_데이터이행_결과서()
    gen_구현단계_완료보고서()

    print("\n=== 테스트 단계 ===")
    gen_통합테스트_계획서()
    gen_통합테스트_시나리오()
    gen_통합테스트_결과서()
    gen_성능테스트_계획서()
    gen_성능테스트_결과서()
    gen_보안점검_결과서()
    gen_UAT_계획서()
    gen_UAT_결과서()
    gen_결함관리_대장()

    print(f"\n=== 완료! docs/ 폴더에 총 32개 산출물 생성 ===")
